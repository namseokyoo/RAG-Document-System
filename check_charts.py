"""
PPT 파일의 차트/그래프 확인
"""
from pptx import Presentation
from pathlib import Path

test_files = [
    "data/test_pptx/advanced_01_financial_report.pptx",
    "data/test_pptx/advanced_02_product_plan.pptx",
    "data/test_pptx/complex_03_data_analysis_report.pptx",
    "data/test_pptx/complex_05_comprehensive_report.pptx",
]

print("=" * 80)
print("PPT 파일별 차트/그래프 현황")
print("=" * 80)
print()

for file_path in test_files:
    test_file = Path(file_path)
    if not test_file.exists():
        continue

    prs = Presentation(str(test_file))

    print(f"파일: {test_file.name}")
    print(f"슬라이드: {len(prs.slides)}개")

    total_charts = 0
    total_images = 0

    for i, slide in enumerate(prs.slides, 1):
        charts = 0
        images = 0

        for shape in slide.shapes:
            if shape.has_chart:
                charts += 1
                total_charts += 1
            elif shape.shape_type == 13:  # PICTURE
                images += 1
                total_images += 1

        if charts > 0 or images > 0:
            print(f"  슬라이드 {i}: 차트 {charts}개, 이미지 {images}개")

    print(f"  총: 차트 {total_charts}개, 이미지 {total_images}개")
    print()
