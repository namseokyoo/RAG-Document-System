"""
PPTX 슬라이드 중심 청킹 엔진
Small-to-Large 아키텍처와 paragraph.level 기반 불릿 그룹핑
Vision-Augmented 청킹 지원
"""
from typing import List, Dict, Any, Optional
from pptx import Presentation
import uuid
import re
import base64
import io
import sys
from .pptx_chunking import PPTXChunk, PPTXChunkMetadata, PPTXChunkFactory, PPTX_CHUNK_TYPE_WEIGHTS
from .chunking_fallback import ChunkingFallback


class PPTXChunkingEngine:
    """PPTX 슬라이드 중심 청킹 엔진"""
    
    def __init__(self, config: Dict[str, Any]):
        self.config = config
        self.fallback = ChunkingFallback(config)
        
        # 설정값들
        self.max_size = config.get("max_size", 300)  # PPTX는 300자 최적
        self.overlap_size = config.get("overlap_size", 50)
        self.enable_small_to_large = config.get("enable_small_to_large", True)
    
    def process_pptx_document(self, pptx_path: str, 
                              enable_vision: bool = False,
                              llm_api_type: str = None,
                              llm_base_url: str = None,
                              llm_model: str = None,
                              llm_api_key: str = None) -> List[PPTXChunk]:
        """PPTX 문서를 슬라이드 구조 인식 기반으로 청킹"""
        all_chunks = []
        document_id = str(uuid.uuid4())
        
        # Vision 청킹을 위해 pptx_path 저장 (COM 렌더링용)
        self.pptx_path = pptx_path if enable_vision else None
        
        try:
            presentation = Presentation(pptx_path)
            
            # Vision 청킹을 위해 presentation 객체 저장 (슬라이드 크기 추출용)
            if enable_vision:
                self.presentation = presentation
            
            # Vision 청킹 사용 시: 모든 슬라이드 이미지를 먼저 렌더링
            slide_images = {}
            if enable_vision and llm_api_type and llm_base_url and llm_model:
                print("[Vision] 모든 슬라이드 이미지 렌더링 중...")
                # COM 방식 사용 시 한 번에 모든 슬라이드 렌더링
                if sys.platform == "win32":
                    try:
                        slide_images = self._render_all_slides_via_com(len(presentation.slides))
                        print(f"[Vision] [OK] COM으로 {len(slide_images)}개 슬라이드 렌더링 완료")
                    except Exception as e:
                        print(f"[Vision] [WARN] COM 방식 실패: {e}, Pillow 방식으로 폴백")
                        # 폴백: 슬라이드별 Pillow 렌더링
                        for slide_index, slide in enumerate(presentation.slides):
                            try:
                                img_base64 = self._slide_to_base64_image(slide, slide_index)

                                # 제목 추출
                                slide_title = ""
                                try:
                                    if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
                                        slide_title = slide.shapes.title.text.strip()
                                except:
                                    pass

                                # 제목이 없으면 슬라이드 번호로 대체
                                if not slide_title:
                                    slide_title = f"제목없음-슬라이드{slide_index + 1}"

                                # 딕셔너리 형식으로 저장
                                slide_images[slide_index] = {
                                    "image": img_base64,
                                    "title": slide_title,
                                    "slide_id": None,  # Pillow는 SlideID 없음
                                    "com_index": slide_index
                                }
                                print(f"  [OK] 슬라이드 {slide_index + 1} 렌더링 완료 (제목: {slide_title})")
                            except Exception as e:
                                print(f"  [WARN] 슬라이드 {slide_index + 1} 렌더링 실패: {e}")
                else:
                    # 비-Windows: Pillow로 슬라이드별 렌더링
                    for slide_index, slide in enumerate(presentation.slides):
                        try:
                            img_base64 = self._slide_to_base64_image(slide, slide_index)

                            # 제목 추출
                            slide_title = ""
                            try:
                                if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
                                    slide_title = slide.shapes.title.text.strip()
                            except:
                                pass

                            # 제목이 없으면 슬라이드 번호로 대체
                            if not slide_title:
                                slide_title = f"제목없음-슬라이드{slide_index + 1}"

                            # 딕셔너리 형식으로 저장
                            slide_images[slide_index] = {
                                "image": img_base64,
                                "title": slide_title,
                                "slide_id": None,  # Pillow는 Sl�라이드ID 없음
                                "com_index": slide_index
                            }
                            print(f"  [OK] 슬라이드 {slide_index + 1} 렌더링 완료 (제목: {slide_title})")
                        except Exception as e:
                            print(f"  [WARN] 슬라이드 {slide_index + 1} 렌더링 실패: {e}")
                
                print(f"[Vision] 총 {len(slide_images)}개 슬라이드 렌더링 완료")
            
            # 슬라이드별 청크 생성
            slides_list = list(presentation.slides)  # Phase 2: 전체 슬라이드 리스트 저장
            for slide_index, slide in enumerate(slides_list):
                slide_number = slide_index + 1
                print(f"슬라이드 {slide_number} 처리 중...")

                # 슬라이드 제목 추출 (메타데이터용) - 제목 없으면 [슬라이드 N] 반환
                slide_title = self._extract_slide_title(slide, slide_num=slide_number)

                # Phase 3: 슬라이드 타입 분류
                slide_type = self._classify_slide_type(slide)
                slide_type_weight = self._get_chunk_weight_by_slide_type(slide_type)
                print(f"  슬라이드 타입: {slide_type} (가중치: {slide_type_weight})")

                # 1. Large 청크 생성 (슬라이드 전체) - Small-to-Large 아키텍처
                if self.enable_small_to_large:
                    if enable_vision and slide_images:
                        # Phase 1-B: 제목 기반 이미지 매칭 (COM과 python-pptx 순서 불일치 해결)
                        matched_image_base64 = self._match_slide_image_by_title(
                            slide_title, slide_index, slide_images
                        )

                        # Phase 1-B: 제목 매칭 실패 시 표 구조 기반 매칭 시도
                        if not matched_image_base64:
                            print(f"  [Vision] 제목 매칭 실패, 표 구조 기반 매칭 시도")
                            matched_image_base64 = self._match_by_table_structure(
                                slide, slide_index, slide_images,
                                llm_api_type, llm_api_key or "", llm_model
                            )

                        if matched_image_base64:
                            # Vision 청킹 사용 (제목 또는 표 구조로 매칭된 이미지 사용)
                            slide_chunk = self._create_slide_summary_chunk_with_vision(
                                slide, document_id, slide_number, slide_title, slide_index,
                                llm_api_type, llm_base_url, llm_model, llm_api_key or "",
                                matched_image_base64,  # 매칭된 이미지 전달
                                slide_type, slide_type_weight  # Phase 3: 슬라이드 타입 전달
                            )
                        else:
                            # 모든 매칭 실패 시 텍스트 청킹으로 폴백
                            print(f"  [WARN] 슬라이드 {slide_number} 이미지 매칭 실패 (제목/표 구조 모두), 텍스트 청킹 사용")
                            slide_chunk = self._create_slide_summary_chunk(
                                slide, document_id, slide_number, slide_title,
                                slides_list, slide_index,
                                slide_type, slide_type_weight
                            )
                    else:
                        # Phase 2: 텍스트 청킹 + 슬라이드 문맥 추가
                        slide_chunk = self._create_slide_summary_chunk(
                            slide, document_id, slide_number, slide_title,
                            slides_list, slide_index,  # Phase 2: 전체 슬라이드와 인덱스 전달
                            slide_type, slide_type_weight  # Phase 3: 슬라이드 타입 전달
                        )
                    all_chunks.append(slide_chunk)
                    parent_id = slide_chunk.id
                else:
                    parent_id = None

                # 2. Small 청크 생성
                slide_chunks = self._process_slide_elements(
                    slide, document_id, slide_number, parent_id, slide_title,
                    slide_type, slide_type_weight  # Phase 3: 슬라이드 타입 전달
                )
                all_chunks.extend(slide_chunks)
        
        except Exception as e:
            print(f"PPTX 처리 중 오류 발생: {e}")
            import traceback
            traceback.print_exc()
            # 에러 발생 시 빈 리스트 반환 (상위에서 폴백 처리)
            return []
        
        print(f"총 {len(all_chunks)}개 청크 생성 완료")
        return all_chunks
    
    def _match_slide_image_by_title(self, slide_title: str, slide_index: int,
                                     slide_images: Dict[int, dict]) -> str:
        """제목 기반 이미지 매칭 (COM과 python-pptx 순서 불일치 해결)

        Args:
            slide_title: python-pptx로 추출한 슬라이드 제목
            slide_index: python-pptx 슬라이드 인덱스 (0-based)
            slide_images: COM으로 렌더링한 이미지 딕셔너리
                         {index: {"image": base64_str, "title": str, ...}}

        Returns:
            매칭된 이미지의 base64 문자열, 또는 None
        """
        # 플레이스홀더 제목 제거 ("제목없음-슬라이드N" → "")
        clean_title = slide_title
        if slide_title.startswith("제목없음"):
            clean_title = ""

        # 1순위: 제목 정확 매칭 (제목이 있는 경우)
        if clean_title:
            for img_data in slide_images.values():
                if img_data["title"] == clean_title:
                    print(f"  [Vision] 제목 매칭 성공: '{clean_title}' (COM index: {img_data['com_index']})")
                    return img_data["image"]

            # 부분 매칭 시도
            for img_data in slide_images.values():
                if clean_title in img_data["title"] or img_data["title"] in clean_title:
                    print(f"  [Vision] 제목 부분 매칭: '{clean_title}' ~ '{img_data['title']}' (COM index: {img_data['com_index']})")
                    return img_data["image"]

        # 2순위: 제목이 없는 슬라이드 - 인덱스 매칭 (폴백)
        # COM과 python-pptx 인덱스가 일치하는 경우만
        if slide_index in slide_images:
            img_data = slide_images[slide_index]
            # 제목이 둘 다 없어야 매칭
            if not clean_title and not img_data["title"]:
                print(f"  [Vision] 인덱스 매칭 (둘 다 제목 없음): {slide_index}")
                return img_data["image"]

        # 3순위: 매칭 실패
        print(f"  [WARN] 이미지 매칭 실패: 제목='{slide_title}', index={slide_index}")
        print(f"  [INFO] 표 구조 기반 매칭을 시도하려면 _match_by_table_structure() 사용")
        return None

    def _extract_slide_title(self, slide, slide_num: int = None) -> str:
        """슬라이드 제목 추출

        Args:
            slide: 분석할 슬라이드
            slide_num: 슬라이드 번호 (제목이 없을 때 대체 텍스트로 사용)

        Returns:
            슬라이드 제목. 제목이 없으면 "제목없음-슬라이드N" 형식 반환
        """
        try:
            if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    return title_text
        except:
            pass

        # 제목이 없을 때 슬라이드 번호로 대체
        if slide_num:
            return f"제목없음-슬라이드{slide_num}"
        return "제목없음"

    def _extract_table_structure(self, slide) -> dict:
        """python-pptx 슬라이드에서 표 구조 추출

        Args:
            slide: python-pptx Slide 객체

        Returns:
            {
                "has_table": bool,
                "table_count": int,
                "tables": [
                    {"rows": int, "cols": int, "headers": [str, ...]},
                    ...
                ]
            }
        """
        result = {
            "has_table": False,
            "table_count": 0,
            "tables": []
        }

        try:
            for shape in slide.shapes:
                if shape.has_table:
                    result["has_table"] = True
                    result["table_count"] += 1

                    table = shape.table
                    rows = len(table.rows)
                    cols = len(table.columns)

                    # 첫 행을 헤더로 추출
                    headers = []
                    if rows > 0:
                        for cell in table.rows[0].cells:
                            try:
                                headers.append(cell.text.strip())
                            except:
                                headers.append("")

                    result["tables"].append({
                        "rows": rows,
                        "cols": cols,
                        "headers": headers
                    })
        except Exception as e:
            print(f"  [WARN] 표 구조 추출 실패: {e}")

        return result

    def _has_chart(self, slide) -> bool:
        """슬라이드에 차트가 있는지 확인

        Args:
            slide: python-pptx Slide 객체

        Returns:
            차트 존재 여부
        """
        try:
            for shape in slide.shapes:
                if shape.has_chart:
                    return True
        except:
            pass
        return False

    def _extract_chart_info(self, slide) -> dict:
        """python-pptx 슬라이드에서 차트 기본 정보 추출

        Note: python-pptx는 차트 타입만 추출 가능, 실제 데이터/트렌드는 Vision 필요

        Args:
            slide: python-pptx Slide 객체

        Returns:
            {
                "has_chart": bool,
                "chart_count": int,
                "charts": [
                    {"type": str, "has_title": bool, "title": str},
                    ...
                ]
            }
        """
        result = {
            "has_chart": False,
            "chart_count": 0,
            "charts": []
        }

        try:
            for shape in slide.shapes:
                if shape.has_chart:
                    result["has_chart"] = True
                    result["chart_count"] += 1

                    chart = shape.chart
                    chart_type = str(chart.chart_type)  # ChartType enum

                    # 차트 제목
                    chart_title = ""
                    try:
                        if chart.has_title and chart.chart_title:
                            chart_title = chart.chart_title.text_frame.text
                    except:
                        pass

                    result["charts"].append({
                        "type": chart_type,
                        "has_title": bool(chart_title),
                        "title": chart_title
                    })
        except Exception as e:
            print(f"  [WARN] 차트 정보 추출 실패: {e}")

        return result

    def _detect_table_structure_via_vision(self, image_base64: str,
                                          llm_api_type: str, llm_api_key: str,
                                          llm_model: str) -> dict:
        """Vision API로 이미지에서 표 구조 감지

        Args:
            image_base64: Base64 인코딩된 슬라이드 이미지
            llm_api_type: "openai" or "ollama"
            llm_api_key: API 키
            llm_model: 모델 이름

        Returns:
            {
                "has_table": bool,
                "table_count": int,
                "tables": [
                    {"rows": int, "cols": int},
                    ...
                ]
            }
        """
        import requests

        result = {
            "has_table": False,
            "table_count": 0,
            "tables": []
        }

        if llm_api_type != "openai":
            print(f"  [WARN] Vision 표 구조 감지는 OpenAI API만 지원")
            return result

        try:
            # 간단한 프롬프트로 표 구조만 감지
            prompt = """이 슬라이드에 표가 있나요? 있다면:

1. 표의 개수
2. 각 표의 행(row) 개수와 열(column) 개수

다음 형식으로만 답변하세요:
표_개수: N
표1: R행 C열
표2: R행 C열
(표가 없으면 "표_개수: 0"만 출력)"""

            endpoint = "https://api.openai.com/v1/chat/completions"
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {llm_api_key}"
            }

            payload = {
                "model": llm_model,
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{image_base64}",
                                    "detail": "high"  # 정확한 표 구조 감지를 위해 high 사용
                                }
                            }
                        ]
                    }
                ],
                "max_tokens": 8000,  # 100 → 8000 (Llama4-scout: 테이블 상세 분석)
                "temperature": 0
            }

            response = requests.post(endpoint, headers=headers, json=payload, timeout=30)
            response.raise_for_status()

            response_text = response.json()["choices"][0]["message"]["content"]

            # 간단한 파싱
            lines = response_text.strip().split('\n')
            for line in lines:
                if "표_개수:" in line or "표 개수:" in line:
                    try:
                        count = int(line.split(':')[1].strip())
                        result["table_count"] = count
                        result["has_table"] = count > 0
                    except:
                        pass
                elif line.startswith("표") and "행" in line and "열" in line:
                    try:
                        # "표1: 3행 3열" 형식 파싱
                        parts = line.split(':')[1].strip().split()
                        rows = int(parts[0].replace('행', ''))
                        cols = int(parts[1].replace('열', ''))
                        result["tables"].append({"rows": rows, "cols": cols})
                    except:
                        pass

        except Exception as e:
            print(f"  [WARN] Vision 표 구조 감지 실패: {e}")

        return result

    def _match_by_table_structure(self, slide, slide_index: int,
                                  slide_images: Dict[int, dict],
                                  llm_api_type: str, llm_api_key: str,
                                  llm_model: str) -> str:
        """표 구조 기반 이미지 매칭 (제목 없는 슬라이드용)

        Args:
            slide: python-pptx Slide 객체
            slide_index: 슬라이드 인덱스
            slide_images: COM으로 렌더링한 이미지 딕셔너리
            llm_api_type: API 타입
            llm_api_key: API 키
            llm_model: 모델 이름

        Returns:
            매칭된 이미지의 base64 문자열, 또는 None
        """
        print(f"  [Vision] 표 구조 기반 매칭 시도...")

        # 1. python-pptx로 실제 슬라이드의 표 구조 추출
        slide_structure = self._extract_table_structure(slide)

        if not slide_structure["has_table"]:
            print(f"  [WARN] 슬라이드에 표 없음, 표 구조 매칭 불가")
            return None

        print(f"  [Vision] 슬라이드 표 구조: {slide_structure['table_count']}개 표")
        for i, tbl in enumerate(slide_structure["tables"], 1):
            print(f"    표{i}: {tbl['rows']}행 {tbl['cols']}열")

        # 2. 각 COM 이미지의 표 구조 감지 및 비교
        best_match = None
        best_score = 0

        for img_idx, img_data in slide_images.items():
            # 실제 제목이 있는 이미지는 건너뜀 (플레이스홀더 제목은 포함)
            img_title = img_data.get("title", "")
            if img_title and not img_title.startswith("제목없음"):
                continue

            print(f"  [Vision] COM 이미지 {img_idx + 1} (제목: {img_title}) 표 구조 감지 중...")

            img_structure = self._detect_table_structure_via_vision(
                img_data["image"], llm_api_type, llm_api_key, llm_model
            )

            if not img_structure["has_table"]:
                print(f"    -> 표 없음")
                continue

            print(f"    -> {img_structure['table_count']}개 표 발견")
            for i, tbl in enumerate(img_structure["tables"], 1):
                print(f"       표{i}: {tbl['rows']}행 {tbl['cols']}열")

            # 3. 표 구조 유사도 계산 (±1 row/col 허용)
            score = 0

            # 표 개수 일치
            if slide_structure["table_count"] == img_structure["table_count"]:
                score += 10

            # 각 표의 행/열 일치도 (±1 허용)
            for s_tbl in slide_structure["tables"]:
                for i_tbl in img_structure["tables"]:
                    row_diff = abs(s_tbl["rows"] - i_tbl["rows"])
                    col_diff = abs(s_tbl["cols"] - i_tbl["cols"])

                    # 완전 일치: +20점
                    if row_diff == 0 and col_diff == 0:
                        score += 20
                        print(f"       표 구조 완전 일치: {s_tbl['rows']}×{s_tbl['cols']}")
                        break
                    # 허용 범위 내 (±1): +15점
                    elif row_diff <= 1 and col_diff <= 1:
                        score += 15
                        print(f"       표 구조 유사 (±1 허용): {s_tbl['rows']}×{s_tbl['cols']} vs {i_tbl['rows']}×{i_tbl['cols']}")
                        break

            print(f"    -> 매칭 점수: {score}")

            if score > best_score:
                best_score = score
                best_match = img_data["image"]

        if best_match and best_score >= 20:  # 최소 1개 표 구조 일치
            print(f"  [Vision] 표 구조 매칭 성공! (점수: {best_score})")
            return best_match
        else:
            print(f"  [WARN] 표 구조 매칭 실패 (최고 점수: {best_score})")
            return None

    def _analyze_chart_with_vision(self, image_base64: str, chart_info: dict,
                                   llm_api_type: str, llm_api_key: str,
                                   llm_model: str, llm_base_url: str = "") -> str:
        """Vision API로 차트 이미지 분석

        Args:
            image_base64: Base64 인코딩된 슬라이드 이미지
            chart_info: _extract_chart_info()로 추출한 차트 정보
            llm_api_type: LLM API 타입
            llm_api_key: LLM API 키
            llm_model: LLM 모델명
            llm_base_url: LLM API Base URL

        Returns:
            차트 분석 결과 텍스트
        """
        # 차트 타입 목록
        chart_types = ", ".join([c["type"] for c in chart_info["charts"]])

        # 프롬프트
        prompt = f"""이 슬라이드에는 {chart_info['chart_count']}개의 차트가 있습니다 (유형: {chart_types}).

각 차트에 대해 다음을 분석하세요:

1. **데이터 유형**: 무엇을 측정하는가? (매출, 성장률, 점유율 등)
2. **주요 트렌드**: 증가/감소/변동 패턴
3. **핵심 인사이트**: 가장 중요한 발견 (최대값, 최소값, 이상치 등)
4. **비교**: 항목 간 비교 (예: A가 B보다 20% 높음)

구조화된 형식으로 답변하세요:
---
차트 1:
- 데이터 유형: ...
- 주요 트렌드: ...
- 핵심 인사이트: ...
- 비교: ...
"""

        # Vision API 호출
        try:
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {llm_api_key}"
            }

            payload = {
                "model": llm_model,
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{image_base64}",
                                    "detail": "high"  # 차트는 고해상도 필요
                                }
                            }
                        ]
                    }
                ],
                "max_tokens": 8000,  # 500 → 8000 (Llama4-scout: 차트/이미지 상세 분석)
                "temperature": 0
            }

            # API URL
            if llm_base_url:
                api_url = f"{llm_base_url}/chat/completions"
            else:
                api_url = "https://api.openai.com/v1/chat/completions"

            response = requests.post(api_url, headers=headers, json=payload, timeout=30)
            response.raise_for_status()

            result = response.json()
            description = result["choices"][0]["message"]["content"]

            return description

        except Exception as e:
            print(f"  [ERROR] 차트 Vision 분석 실패: {e}")
            return f"[차트 분석 실패: {e}]"

    def _create_slide_summary_chunk(self, slide, document_id: str, slide_num: int,
                                   slide_title: str, slides_list: List = None,
                                   slide_index: int = None,
                                   slide_type: str = None, slide_type_weight: float = 1.0) -> PPTXChunk:
        """슬라이드 전체 텍스트를 부모 청크로 생성 - Phase 2 & 3 Enhanced"""
        slide_text = self._extract_full_text_from_slide(slide)

        # Phase 2: 슬라이드 문맥 추가 (이전/다음 슬라이드 제목)
        if slides_list is not None and slide_index is not None:
            slide_text = self._add_slide_context(slides_list, slide_index, slide_text)

        # Phase 3: 슬라이드 타입 정보 추가
        chunk = PPTXChunkFactory.create_slide_summary_chunk(
            slide_text=slide_text,
            document_id=document_id,
            slide_num=slide_num,
            slide_title=slide_title
        )

        # Phase 3: 슬라이드 타입별 가중치 적용
        if slide_type:
            chunk.metadata.slide_type = slide_type
            chunk.metadata.chunk_type_weight = chunk.metadata.chunk_type_weight * slide_type_weight

        return chunk
    
    def _extract_full_text_from_slide(self, slide) -> str:
        """슬라이드 내의 모든 보이는 텍스트 추출 (노트 제외)"""
        full_text = []
        
        try:
            # 1. 제목
            if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    full_text.append(f"[제목]:\n{title_text}")
            
            # 2. 본문 및 기타 텍스트
            for shape in slide.shapes:
                # 제목 중복 방지
                if shape == slide.shapes.title:
                    continue

                # 텍스트 프레임 - 개별 try-except로 보호
                try:
                    if hasattr(shape, 'text_frame') and shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            full_text.append(f"[본문]:\n{text}")
                except Exception:
                    # 텍스트 추출 실패 시 조용히 건너뜀 (전체 텍스트 추출 컨텍스트)
                    pass

                # 테이블 - 개별 try-except로 보호
                try:
                    if hasattr(shape, 'table') and shape.has_table:
                        table_text = self._convert_table_to_simple_text(shape.table)
                        if table_text:
                            full_text.append(f"[표]:\n{table_text}")
                except Exception:
                    # 테이블 추출 실패 시 조용히 건너뜀 (전체 텍스트 추출 컨텍스트)
                    pass
        
        except Exception as e:
            print(f"슬라이드 텍스트 추출 중 오류: {e}")
        
        return "\n\n".join(full_text) if full_text else ""
    
    def _process_slide_elements(self, slide, document_id: str, slide_num: int,
                               parent_id: Optional[str], slide_title: str,
                               slide_type: str = None, slide_type_weight: float = 1.0) -> List[PPTXChunk]:
        """슬라이드 요소들을 Small 청크로 처리 - Phase 3 Enhanced"""
        chunks = []
        table_idx = 0  # 슬라이드별 표 인덱스 추적

        try:
            # 2a. 슬라이드 제목
            title_chunk = self._create_slide_title_chunk(
                slide, document_id, slide_num, parent_id, slide_title
            )
            if title_chunk:
                # Phase 3: 슬라이드 타입별 가중치 적용
                if slide_type:
                    title_chunk.metadata.slide_type = slide_type
                    title_chunk.metadata.chunk_type_weight = title_chunk.metadata.chunk_type_weight * slide_type_weight
                chunks.append(title_chunk)
            
            # 2b. 슬라이드 노트
            notes_chunks = self._create_slide_notes_chunks(
                slide, document_id, slide_num, parent_id, slide_title
            )
            # Phase 3: 슬라이드 타입별 가중치 적용
            for chunk in notes_chunks:
                if slide_type:
                    chunk.metadata.slide_type = slide_type
                    chunk.metadata.chunk_type_weight = chunk.metadata.chunk_type_weight * slide_type_weight
            chunks.extend(notes_chunks)

            # 2c. 슬라이드 본문 (도형 순회)
            shape_index = 0
            for shape in slide.shapes:
                # 제목 중복 방지
                if shape == slide.shapes.title:
                    continue

                shape_index += 1

                # 텍스트 프레임 (불릿 포인트) - 개별 try-except로 보호
                try:
                    if hasattr(shape, 'text_frame') and shape.has_text_frame:
                        bullet_chunks = self._chunk_bullet_points_by_level(
                            shape.text_frame, document_id, slide_num, parent_id, slide_title
                        )
                        # 메타데이터에 shape 정보 추가 (가능한 경우)
                        for chunk in bullet_chunks:
                            if hasattr(shape, 'shape_type'):
                                chunk.metadata.shape_type = str(shape.shape_type)
                            # Phase 3: 슬라이드 타입별 가중치 적용
                            if slide_type:
                                chunk.metadata.slide_type = slide_type
                                chunk.metadata.chunk_type_weight = chunk.metadata.chunk_type_weight * slide_type_weight
                        chunks.extend(bullet_chunks)
                except Exception as e:
                    # shape.has_text_frame 체크 후에도 실제 접근 시 오류 발생 가능
                    print(f"  [WARN] Shape {shape_index} 텍스트 처리 중 오류 (건너뜀): {e}")

                # 테이블 - 개별 try-except로 보호
                try:
                    if hasattr(shape, 'table') and shape.has_table:
                        table_chunks = self._chunk_pptx_table(
                            shape.table, document_id, slide_num, parent_id, slide_title, table_idx
                        )
                        # Phase 3: 슬라이드 타입별 가중치 적용
                        for chunk in table_chunks:
                            if slide_type:
                                chunk.metadata.slide_type = slide_type
                                chunk.metadata.chunk_type_weight = chunk.metadata.chunk_type_weight * slide_type_weight
                        chunks.extend(table_chunks)
                        table_idx += 1  # 다음 표를 위해 인덱스 증가
                except Exception as e:
                    # shape.has_table 체크 후에도 shape.table 접근 시 오류 발생 가능
                    print(f"  [WARN] Shape {shape_index} 테이블 처리 중 오류 (건너뜀): {e}")
        
        except Exception as e:
            print(f"슬라이드 요소 처리 중 오류: {e}")
            import traceback
            traceback.print_exc()
        
        return chunks
    
    def _create_slide_title_chunk(self, slide, document_id: str, slide_num: int,
                                 parent_id: Optional[str], slide_title: str) -> Optional[PPTXChunk]:
        """슬라이드 제목 청크 생성"""
        try:
            if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    return PPTXChunkFactory.create_slide_title_chunk(
                        title_text=title_text,
                        document_id=document_id,
                        slide_num=slide_num,
                        parent_id=parent_id,
                        slide_title=slide_title
                    )
        except Exception as e:
            print(f"제목 청크 생성 중 오류: {e}")
        
        return None
    
    def _create_slide_notes_chunks(self, slide, document_id: str, slide_num: int,
                                  parent_id: Optional[str], slide_title: str) -> List[PPTXChunk]:
        """슬라이드 노트 청크 생성"""
        chunks = []
        
        try:
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes_slide = slide.notes_slide
                if hasattr(notes_slide, 'notes_text_frame'):
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        # 노트가 길 수 있으므로 Fallback 적용
                        from .pdf_chunking import ChunkMetadata
                        
                        base_metadata = ChunkMetadata(
                            document_id=document_id,
                            page_number=slide_num,
                            parent_chunk_id=parent_id,
                            section_title=slide_title,
                            chunk_type_weight=PPTX_CHUNK_TYPE_WEIGHTS.get("slide_notes", 1.7)
                        )
                        
                        fallback_chunks = self.fallback.chunk_element_with_fallback(
                            content=notes_text,
                            element_type="slide_notes",
                            base_metadata=base_metadata
                        )
                        
                        # PDF Chunk를 PPTX Chunk로 변환
                        for chunk in fallback_chunks:
                            pptx_metadata = PPTXChunkMetadata(
                                document_id=document_id,
                                slide_number=slide_num,
                                parent_chunk_id=parent_id,
                                slide_title=slide_title,
                                chunk_type_weight=1.7,
                                has_notes=True
                            )
                            
                            pptx_chunk = PPTXChunk(
                                id=chunk.id,
                                content=chunk.content,
                                chunk_type=chunk.chunk_type,
                                metadata=pptx_metadata
                            )
                            chunks.append(pptx_chunk)
        
        except Exception as e:
            print(f"노트 청크 생성 중 오류: {e}")
        
        return chunks
    
    def _chunk_bullet_points_by_level(self, text_frame, document_id: str, slide_num: int,
                                     parent_id: Optional[str], slide_title: str) -> List[PPTXChunk]:
        """paragraph.level을 사용하여 불릿을 의미론적 그룹으로 묶기"""
        chunks = []
        current_bullet_group_lines = []
        current_group_top_level = None  # 그룹의 최상위 레벨 추적
        current_group_levels = []  # 그룹 내 사용된 모든 레벨 추적
        
        try:
            # TextFrame 내의 모든 문단(불릿) 순회
            for paragraph in text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if not para_text:
                    continue
                
                # 들여쓰기 레벨 추출
                indent_level = paragraph.level
                formatted_text = f"{'  ' * indent_level}{para_text}"
                
                # 새 최상위 불릿(level 0)이 시작되고, 버퍼에 내용이 있다면
                if indent_level == 0 and current_bullet_group_lines:
                    # 지금까지 모은 버퍼를 하나의 청크로
                    group_content = "\n".join(current_bullet_group_lines)
                    
                    # 그룹의 최상위 레벨과 최대 레벨 계산
                    group_top_level = current_group_top_level if current_group_top_level is not None else 0
                    group_max_level = max(current_group_levels) if current_group_levels else 0
                    
                    # Fallback 적용 (너무 크면 분할)
                    from .pdf_chunking import ChunkMetadata
                    
                    base_metadata = ChunkMetadata(
                        document_id=document_id,
                        page_number=slide_num,
                        parent_chunk_id=parent_id,
                        section_title=slide_title,
                        chunk_type_weight=PPTX_CHUNK_TYPE_WEIGHTS.get("bullet_group", 1.0)
                    )
                    
                    fallback_chunks = self.fallback.chunk_element_with_fallback(
                        content=group_content,
                        element_type="bullet_group",
                        base_metadata=base_metadata
                    )
                    
                    # PDF Chunk를 PPTX Chunk로 변환
                    for chunk in fallback_chunks:
                        pptx_metadata = PPTXChunkMetadata(
                            document_id=document_id,
                            slide_number=slide_num,
                            parent_chunk_id=parent_id,
                            slide_title=slide_title,
                            chunk_type_weight=1.0,
                            bullet_level=group_top_level  # 그룹의 최상위 레벨 저장
                        )
                        
                        pptx_chunk = PPTXChunk(
                            id=chunk.id,
                            content=chunk.content,
                            chunk_type=chunk.chunk_type,
                            metadata=pptx_metadata
                        )
                        chunks.append(pptx_chunk)
                    
                    # 버퍼 비우고 새 최상위 불릿으로 시작
                    current_bullet_group_lines = [formatted_text]
                    current_group_top_level = indent_level
                    current_group_levels = [indent_level]
                else:
                    # 하위 불릿이거나 첫 번째 상위 불릿이면 버퍼에 추가
                    current_bullet_group_lines.append(formatted_text)
                    if current_group_top_level is None:
                        current_group_top_level = indent_level
                    current_group_levels.append(indent_level)
            
            # 루프 종료 후, 마지막 불릿 그룹 처리
            if current_bullet_group_lines:
                group_content = "\n".join(current_bullet_group_lines)
                
                # 그룹의 최상위 레벨 계산
                group_top_level = current_group_top_level if current_group_top_level is not None else 0
                
                from .pdf_chunking import ChunkMetadata
                
                base_metadata = ChunkMetadata(
                    document_id=document_id,
                    page_number=slide_num,
                    parent_chunk_id=parent_id,
                    section_title=slide_title,
                    chunk_type_weight=PPTX_CHUNK_TYPE_WEIGHTS.get("bullet_group", 1.0)
                )
                
                fallback_chunks = self.fallback.chunk_element_with_fallback(
                    content=group_content,
                    element_type="bullet_group",
                    base_metadata=base_metadata
                )
                
                for chunk in fallback_chunks:
                    pptx_metadata = PPTXChunkMetadata(
                        document_id=document_id,
                        slide_number=slide_num,
                        parent_chunk_id=parent_id,
                        slide_title=slide_title,
                        chunk_type_weight=1.0,
                        bullet_level=group_top_level  # 그룹의 최상위 레벨 저장
                    )
                    
                    pptx_chunk = PPTXChunk(
                        id=chunk.id,
                        content=chunk.content,
                        chunk_type=chunk.chunk_type,
                        metadata=pptx_metadata
                    )
                    chunks.append(pptx_chunk)
        
        except Exception as e:
            print(f"불릿 포인트 청킹 중 오류: {e}")
        
        return chunks
    
    def _chunk_pptx_table(self, table, document_id: str, slide_num: int,
                         parent_id: Optional[str], slide_title: str, table_idx: int = 0) -> List[PPTXChunk]:
        """Phase 1-3: 표 다층 청킹 - 전체/행/열/키-값 청크 생성"""
        chunks = []
        
        try:
            if len(table.rows) == 0:
                return chunks
            
            # 표 기본 정보 추출 (슬라이드별 표 인덱스 사용)
            table_id = f"{document_id}_slide_{slide_num}_table_{table_idx}"
            num_rows = len(table.rows)
            num_cols = len(table.rows[0].cells) if num_rows > 0 else 0
            
            # 헤더 행 추출
            header_row = []
            if num_rows > 0:
                header_row = [cell.text.strip() for cell in table.rows[0].cells]
            
            # 표 제목 추출 (슬라이드 제목 또는 첫 번째 셀에서)
            table_title = slide_title or header_row[0] if header_row else None
            
            # 항목 번호 추출 (Phase 2)
            item_map = self._extract_item_numbers_from_table(table, header_row)
            
            # 데이터 타입 감지 (Phase 3)
            data_type = self._detect_table_data_type(header_row, table)
            
            # 1. 전체 표 청크 생성 (컨텍스트용)
            full_table_chunk = self._create_full_table_chunk(
                table, document_id, slide_num, parent_id, slide_title,
                table_id, table_title, header_row, num_rows, num_cols, data_type
            )
            if full_table_chunk:  # None이 아닐 때만 추가
                chunks.append(full_table_chunk)
            
            # 2. 각 행을 개별 청크로 생성
            for row_idx, row in enumerate(table.rows):
                row_chunk = self._create_table_row_chunk(
                    row, row_idx, header_row, document_id, slide_num, parent_id,
                    slide_title, table_id, table_title, num_cols, item_map, data_type
                )
                chunks.append(row_chunk)
            
            # 3. 각 열을 개별 청크로 생성 (열별 집계 검색용)
            if num_cols > 0:
                for col_idx in range(num_cols):
                    col_chunk = self._create_table_column_chunk(
                        table, col_idx, header_row, document_id, slide_num, parent_id,
                        slide_title, table_id, table_title, num_rows, data_type
                    )
                    if col_chunk:  # None이 아닐 때만 추가
                        chunks.append(col_chunk)
            
            # 4. 키-값 쌍 청크 생성 (항목 번호 검색용, Phase 2)
            kv_chunks = self._create_table_key_value_chunks(
                table, header_row, document_id, slide_num, parent_id,
                slide_title, table_id, table_title, item_map, data_type
            )
            chunks.extend(kv_chunks)
        
        except Exception as e:
            print(f"테이블 청킹 중 오류: {e}")
            import traceback
            traceback.print_exc()
        
        return chunks
    
    def _create_full_table_chunk(self, table, document_id: str, slide_num: int,
                                parent_id: Optional[str], slide_title: str,
                                table_id: str, table_title: Optional[str],
                                header_row: List[str], num_rows: int, num_cols: int,
                                data_type: Optional[str]) -> Optional[PPTXChunk]:
        """전체 표 청크 생성 (컨텍스트용) - Phase 1 Enhanced"""
        markdown_table = self._convert_table_to_markdown(table)

        # 빈 표는 생성하지 않음
        if not markdown_table or not markdown_table.strip():
            return None

        # Phase 1: 표 요약 및 숫자 메타데이터 추가
        table_summary = self._generate_table_summary(table, header_row, table_title)
        numbers = self._extract_numbers_from_table(table)

        # Phase 1: 표 내용에 요약 포함
        enhanced_content = f"""[표 요약]
{table_summary}

[표 데이터]
{markdown_table}"""

        if numbers:
            # 숫자가 있으면 메타정보 추가
            enhanced_content += f"\n\n[주요 숫자] {', '.join(numbers[:10])}"  # 최대 10개만

        metadata = PPTXChunkMetadata(
            document_id=document_id,
            slide_number=slide_num,
            parent_chunk_id=parent_id,
            slide_title=slide_title,
            chunk_type_weight=1.2,  # 전체 표는 높은 가중치 (집계 질문에 중요)
            has_table=True,
            shape_type="table",
            table_id=table_id,
            table_title=table_title,
            header_row=header_row,
            table_row_count=num_rows,
            table_col_count=num_cols,
            data_type=data_type
        )

        return PPTXChunk(
            id=f"{table_id}_full",
            content=enhanced_content,
            chunk_type="table_full",
            metadata=metadata
        )
    
    def _create_table_row_chunk(self, row, row_idx: int, header_row: List[str],
                                document_id: str, slide_num: int,
                                parent_id: Optional[str], slide_title: str,
                                table_id: str, table_title: Optional[str],
                                num_cols: int, item_map: Dict[str, Dict],
                                data_type: Optional[str]) -> PPTXChunk:
        """행 단위 청크 생성"""
        cells = [cell.text.strip() for cell in row.cells]
        is_header = (row_idx == 0)
        
        # 행 데이터를 텍스트로 변환
        if is_header:
            row_text = f"헤더: {' | '.join(cells)}"
        else:
            # 헤더와 함께 키-값 쌍 형식으로 변환 (더 명확한 형식)
            row_pairs = []
            for col_idx, cell_text in enumerate(cells):
                if col_idx < len(header_row):
                    header_name = header_row[col_idx]
                    row_pairs.append(f"{header_name}: {cell_text}")
                else:
                    # 헤더가 없는 열도 포함
                    row_pairs.append(f"열{col_idx + 1}: {cell_text}")
            row_text = " | ".join(row_pairs)
        
        # 항목 번호 추출 (Phase 2) - 모든 셀에서 검색
        item_number = None
        if cells and not is_header:
            # 첫 번째 셀뿐만 아니라 모든 셀에서 검색
            for cell_text in cells:
                if match := re.search(r'항목\s*(\d+)', cell_text):
                    item_number = f"항목 {match.group(1)}"
                    break  # 첫 번째 매치만 사용
        
        # 셀 참조 생성 (Phase 3) - 모든 셀 참조 생성
        if not is_header and cells:
            # 첫 번째 셀의 참조만 저장 (전체 행은 모든 셀 포함)
            cell_reference = f"R{row_idx + 1}C1"
            # 전체 행의 셀 참조도 생성 가능 (향후 확장용)
        else:
            cell_reference = None
        
        metadata = PPTXChunkMetadata(
            document_id=document_id,
            slide_number=slide_num,
            parent_chunk_id=parent_id,
            slide_title=slide_title,
            chunk_type_weight=1.3,  # 행 단위는 높은 가중치
            has_table=True,
            shape_type="table",
            table_id=table_id,
            table_title=table_title,
            row_index=row_idx,
            cell_reference=cell_reference,
            header_row=header_row,
            is_header_row=is_header,
            item_number=item_number,
            data_type=data_type,
            table_col_count=num_cols
        )
        
        return PPTXChunk(
            id=f"{table_id}_row_{row_idx}",
            content=row_text,
            chunk_type="table_row",
            metadata=metadata
        )
    
    def _create_table_column_chunk(self, table, col_idx: int, header_row: List[str],
                                   document_id: str, slide_num: int,
                                   parent_id: Optional[str], slide_title: str,
                                   table_id: str, table_title: Optional[str],
                                   num_rows: int, data_type: Optional[str]) -> Optional[PPTXChunk]:
        """열 단위 청크 생성 (열별 집계 검색용)"""
        if col_idx >= len(header_row):
            return None
        
        col_header = header_row[col_idx]
        col_values = []
        
        # 해당 열의 모든 값 추출
        rows_list = list(table.rows)  # 슬라이싱을 위해 리스트로 변환
        for row_idx, row in enumerate(rows_list[1:], start=1):  # 헤더 제외
            if col_idx < len(row.cells):
                cell_text = row.cells[col_idx].text.strip()
                col_values.append(cell_text)
        
        # 열 데이터를 텍스트로 변환
        col_text = f"{col_header} 열: {', '.join(col_values)}"
        
        metadata = PPTXChunkMetadata(
            document_id=document_id,
            slide_number=slide_num,
            parent_chunk_id=parent_id,
            slide_title=slide_title,
            chunk_type_weight=1.1,  # 열 단위는 중간 가중치
            has_table=True,
            shape_type="table",
            table_id=table_id,
            table_title=table_title,
            col_index=col_idx,
            header_row=header_row,
            data_type=data_type,
            table_row_count=num_rows
        )
        
        return PPTXChunk(
            id=f"{table_id}_col_{col_idx}",
            content=col_text,
            chunk_type="table_column",
            metadata=metadata
        )
    
    def _create_table_key_value_chunks(self, table, header_row: List[str],
                                       document_id: str, slide_num: int,
                                       parent_id: Optional[str], slide_title: str,
                                       table_id: str, table_title: Optional[str],
                                       item_map: Dict[str, Dict],
                                       data_type: Optional[str]) -> List[PPTXChunk]:
        """키-값 쌍 청크 생성 (항목 번호 검색용, Phase 2)"""
        chunks = []
        
        for item_number, item_info in item_map.items():
            row_idx = item_info["row_index"]
            row_data = item_info["full_row_data"]
            
            # 키-값 쌍 형식으로 변환
            kv_pairs = []
            for col_idx, value in enumerate(row_data):
                if col_idx < len(header_row):
                    kv_pairs.append(f"{header_row[col_idx]}: {value}")
            
            kv_text = f"{item_number} - {' | '.join(kv_pairs)}"
            
            metadata = PPTXChunkMetadata(
                document_id=document_id,
                slide_number=slide_num,
                parent_chunk_id=parent_id,
                slide_title=slide_title,
                chunk_type_weight=1.5,  # 항목 번호는 높은 가중치
                has_table=True,
                shape_type="table",
                table_id=table_id,
                table_title=table_title,
                row_index=row_idx,
                cell_reference=f"R{row_idx + 1}C1",
                header_row=header_row,
                item_number=item_number,
                data_type=data_type
            )
            
            chunk = PPTXChunk(
                id=f"{table_id}_item_{item_number}",
                content=kv_text,
                chunk_type="table_key_value",
                metadata=metadata
            )
            chunks.append(chunk)
        
        return chunks
    
    def _extract_item_numbers_from_table(self, table, header_row: List[str]) -> Dict[str, Dict]:
        """표에서 항목 번호 추출 (Phase 2)"""
        item_map = {}
        
        rows_list = list(table.rows)  # 슬라이싱을 위해 리스트로 변환
        if len(rows_list) <= 1:
            return item_map
        
        # 첫 번째 열에서 항목 번호 추출
        for row_idx, row in enumerate(rows_list[1:], start=1):  # 헤더 제외
            if len(row.cells) > 0:
                first_cell = row.cells[0].text.strip()
                
                # "항목 1", "항목 2" 패턴 추출
                if match := re.search(r'항목\s*(\d+)', first_cell):
                    item_num = match.group(1)
                    item_number = f"항목 {item_num}"
                    
                    # 전체 행 데이터 추출
                    full_row_data = [cell.text.strip() for cell in row.cells]
                    
                    item_map[item_number] = {
                        "row_index": row_idx,
                        "full_row_data": full_row_data
                    }
        
        return item_map
    
    def _detect_table_data_type(self, header_row: List[str], table) -> Optional[str]:
        """표 데이터 타입 자동 감지 (Phase 3)"""
        if not header_row:
            return None
        
        header_text = " ".join(header_row).lower()
        
        # 예산 관련 키워드
        if any(keyword in header_text for keyword in ["예산", "budget", "지출", "비용"]):
            return "budget"
        
        # 매출 관련 키워드
        if any(keyword in header_text for keyword in ["매출", "sales", "수익", "revenue"]):
            return "sales"
        
        # 성과 관련 키워드
        if any(keyword in header_text for keyword in ["성과", "performance", "점수", "score"]):
            return "performance"
        
        # 일정 관련 키워드
        if any(keyword in header_text for keyword in ["일정", "schedule", "기간", "period"]):
            return "schedule"
        
        return None
    
    def _convert_table_to_markdown(self, table) -> str:
        """테이블을 Markdown 형식으로 변환"""
        markdown_lines = []
        
        try:
            rows_list = list(table.rows)  # 슬라이싱을 위해 리스트로 변환
            if len(rows_list) == 0:
                return ""
            
            # 헤더 (첫 번째 행)
            header_cells = [cell.text.strip() for cell in rows_list[0].cells]
            markdown_lines.append(f"| {' | '.join(header_cells)} |")
            markdown_lines.append(f"| {' | '.join(['---'] * len(header_cells))} |")
            
            # 본문 (나머지 행)
            for row in rows_list[1:]:
                body_cells = [cell.text.strip() for cell in row.cells]
                markdown_lines.append(f"| {' | '.join(body_cells)} |")
            
            return "\n".join(markdown_lines)
        
        except Exception as e:
            print(f"테이블 변환 중 오류: {e}")
            import traceback
            traceback.print_exc()
            return ""
    
    def _convert_table_to_simple_text(self, table) -> str:
        """테이블을 단순 텍스트로 변환"""
        text_lines = []

        try:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                if row_text:
                    text_lines.append(row_text)

            return "\n".join(text_lines)

        except Exception as e:
            print(f"테이블 텍스트 변환 중 오류: {e}")
            return ""

    # ========== Phase 1: Table Structure Enhancement ==========

    def _extract_numbers_from_table(self, table) -> List[str]:
        """표에서 모든 숫자 추출 (숫자 검색 최적화)"""
        numbers = []

        try:
            for row in table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    # 숫자 패턴 추출 (쉼표 포함, 소수점, 백분율, 통화 등)
                    number_matches = re.findall(r'[\d,]+\.?\d*%?|[₩$€£¥]\s*[\d,]+\.?\d*', text)
                    numbers.extend(number_matches)

        except Exception as e:
            print(f"표 숫자 추출 중 오류: {e}")

        return numbers

    def _generate_table_summary(self, table, header_row: List[str], table_title: str = None) -> str:
        """표의 자연어 요약 생성"""
        try:
            num_rows = len(list(table.rows))
            num_cols = len(header_row) if header_row else 0

            # 기본 구조 설명
            summary_parts = []

            if table_title:
                summary_parts.append(f"표 제목: {table_title}")

            summary_parts.append(f"{num_rows}행 x {num_cols}열 표")

            # 헤더 정보
            if header_row and len(header_row) > 0:
                header_desc = ", ".join([f"'{h}'" for h in header_row[:5]])  # 최대 5개만
                if len(header_row) > 5:
                    header_desc += f" 외 {len(header_row) - 5}개"
                summary_parts.append(f"열: {header_desc}")

            # 숫자 정보
            numbers = self._extract_numbers_from_table(table)
            if numbers:
                summary_parts.append(f"숫자 데이터 {len(numbers)}개 포함")

            return " | ".join(summary_parts)

        except Exception as e:
            print(f"표 요약 생성 중 오류: {e}")
            return "표 데이터"

    # ========== Phase 2: Slide Context Enhancement ==========

    def _get_slide_title(self, slide) -> str:
        """슬라이드 제목 추출"""
        try:
            if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
                title = slide.shapes.title.text.strip()
                if title:
                    return title
        except:
            pass
        return "제목 없음"

    def _add_slide_context(self, slides: List, slide_index: int, current_content: str) -> str:
        """이전/다음 슬라이드 제목을 컨텍스트로 추가"""
        context_parts = []

        # 이전 슬라이드 제목
        if slide_index > 0:
            prev_title = self._get_slide_title(slides[slide_index - 1])
            context_parts.append(f"[이전 슬라이드] {prev_title}")

        # 현재 슬라이드 마커
        context_parts.append("[현재 슬라이드]")
        context_parts.append(current_content)

        # 다음 슬라이드 제목
        if slide_index < len(slides) - 1:
            next_title = self._get_slide_title(slides[slide_index + 1])
            context_parts.append(f"[다음 슬라이드] {next_title}")

        return "\n\n".join(context_parts)

    # ========== Phase 3: Slide Type Classification ==========

    def _classify_slide_type(self, slide) -> str:
        """슬라이드 타입 자동 분류"""
        try:
            table_count = 0
            chart_count = 0
            image_count = 0
            text_shape_count = 0
            total_text_length = 0
            bullet_count = 0

            # 슬라이드 요소 분석
            for shape in slide.shapes:
                # 표 카운트
                if hasattr(shape, 'has_table'):
                    try:
                        if shape.has_table:
                            table_count += 1
                    except:
                        pass

                # 차트 카운트
                if hasattr(shape, 'has_chart'):
                    try:
                        if shape.has_chart:
                            chart_count += 1
                    except:
                        pass

                # 이미지 카운트 (차트 제외)
                if hasattr(shape, 'shape_type'):
                    try:
                        # shape_type 13 = PICTURE
                        if str(shape.shape_type) == '13':
                            image_count += 1
                    except:
                        pass

                # 텍스트 분석
                if hasattr(shape, 'text'):
                    try:
                        text = shape.text.strip()
                        if text:
                            text_shape_count += 1
                            total_text_length += len(text)

                            # 불릿 포인트 카운트
                            if hasattr(shape, 'text_frame'):
                                for paragraph in shape.text_frame.paragraphs:
                                    if hasattr(paragraph, 'level') and paragraph.text.strip():
                                        bullet_count += 1
                    except:
                        pass

            # 타입 결정 로직
            if table_count >= 2:
                return "table_heavy"  # 표 2개 이상
            elif table_count == 1 and total_text_length < 200:
                return "table_focused"  # 표 1개 + 적은 텍스트
            elif chart_count >= 2:
                return "chart_heavy"  # 차트 2개 이상
            elif chart_count == 1 and total_text_length < 200:
                return "chart_focused"  # 차트 1개 + 적은 텍스트
            elif bullet_count >= 5:
                return "bullet_list"  # 불릿 포인트 5개 이상
            elif total_text_length > 500:
                return "text_heavy"  # 긴 텍스트
            elif image_count >= 2:
                return "image_heavy"  # 이미지 2개 이상
            elif total_text_length < 100:
                return "minimal"  # 최소 텍스트 (제목만 등)
            else:
                return "mixed"  # 혼합형

        except Exception as e:
            print(f"슬라이드 타입 분류 오류: {e}")
            return "unknown"

    def _get_chunk_weight_by_slide_type(self, slide_type: str) -> float:
        """슬라이드 타입별 청크 가중치 반환"""
        weights = {
            "table_heavy": 1.3,      # 표 중심 슬라이드 - 높은 가중치
            "table_focused": 1.3,
            "chart_heavy": 1.2,      # 차트 중심 슬라이드
            "chart_focused": 1.2,
            "bullet_list": 1.0,      # 불릿 리스트 - 기본 가중치
            "text_heavy": 0.9,       # 텍스트 중심 - 낮은 가중치
            "image_heavy": 1.1,      # 이미지 중심
            "minimal": 0.7,          # 최소 내용 - 낮은 가중치
            "mixed": 1.0,            # 혼합형 - 기본 가중치
            "unknown": 1.0
        }
        return weights.get(slide_type, 1.0)

    def _get_optimal_chunk_size_by_type(self, slide_type: str) -> int:
        """슬라이드 타입별 최적 청크 크기 반환 (Phase 6 준비)"""
        sizes = {
            "table_heavy": 500,      # 표는 크게
            "table_focused": 500,
            "chart_heavy": 400,      # 차트는 중간
            "chart_focused": 400,
            "bullet_list": 300,      # 불릿은 기본
            "text_heavy": 400,       # 긴 텍스트는 중간
            "image_heavy": 300,      # 이미지는 기본
            "minimal": 200,          # 최소 내용은 작게
            "mixed": 300,
            "unknown": 300
        }
        return sizes.get(slide_type, self.max_size)
    
    def get_chunk_statistics(self, chunks: List[PPTXChunk]) -> Dict[str, Any]:
        """청크 통계 정보 반환"""
        if not chunks:
            return {}
        
        stats = {
            "total_chunks": len(chunks),
            "chunk_types": {},
            "avg_word_count": 0,
            "avg_char_count": 0,
            "slides_covered": set(),
            "slides_with_notes": 0
        }
        
        total_words = 0
        total_chars = 0
        
        for chunk in chunks:
            # 청크 타입별 카운트
            chunk_type = chunk.chunk_type
            stats["chunk_types"][chunk_type] = stats["chunk_types"].get(chunk_type, 0) + 1
            
            # 단어/문자 수 누적
            total_words += chunk.metadata.word_count
            total_chars += chunk.metadata.char_count
            
            # 슬라이드 정보
            stats["slides_covered"].add(chunk.metadata.slide_number)
            if chunk.metadata.has_notes:
                stats["slides_with_notes"] += 1
        
        stats["avg_word_count"] = total_words / len(chunks)
        stats["avg_char_count"] = total_chars / len(chunks)
        stats["slides_covered"] = len(stats["slides_covered"])
        
        return stats
    
    # ===== Vision-Augmented Chunking Methods =====
    
    def _create_slide_summary_chunk_with_vision(self, slide, document_id: str,
                                                slide_num: int, slide_title: str, slide_index: int,
                                                llm_api_type: str, llm_base_url: str,
                                                llm_model: str, llm_api_key: str,
                                                slide_img_base64: str = None,
                                                slide_type: str = None, slide_type_weight: float = 1.0) -> PPTXChunk:
        """슬라이드 전체를 부모 청크로 생성 (Vision LLM 사용) - Phase 3 Enhanced"""

        # 1. 기본 텍스트 추출
        slide_text = self._extract_full_text_from_slide(slide)

        # 2. Vision LLM으로 시각적 분석 추가 (이미 렌더링된 이미지 사용)
        try:
            # total_slides 계산 (presentation 객체에서)
            total_slides = len(self.presentation.slides) if hasattr(self, 'presentation') else 0

            vision_description = self._analyze_slide_with_vision(
                slide, slide_index, llm_api_type, llm_base_url, llm_model, llm_api_key,
                slide_img_base64,  # 미리 렌더링된 이미지 전달
                slide_title=slide_title,
                slide_num=slide_num,
                total_slides=total_slides
            )
            # Vision 설명을 텍스트 앞에 추가
            enhanced_text = f"""[Vision Analysis]
{vision_description}

[Original Content]
{slide_text}"""
        except Exception as e:
            print(f"[WARN] Vision 분석 실패 (슬라이드 {slide_num}): {e}, 텍스트만 사용")
            enhanced_text = slide_text

        # Phase 3: 슬라이드 타입 정보 추가
        chunk = PPTXChunkFactory.create_slide_summary_chunk(
            slide_text=enhanced_text,
            document_id=document_id,
            slide_num=slide_num,
            slide_title=slide_title
        )

        # Phase 3: 슬라이드 타입별 가중치 적용
        if slide_type:
            chunk.metadata.slide_type = slide_type
            chunk.metadata.chunk_type_weight = chunk.metadata.chunk_type_weight * slide_type_weight

        return chunk
    
    def _analyze_slide_with_vision(self, slide, slide_index: int, llm_api_type: str,
                                   llm_base_url: str, llm_model: str,
                                   llm_api_key: str, slide_img_base64: str = None,
                                   slide_title: str = "", slide_num: int = 0, total_slides: int = 0) -> str:
        """Vision LLM으로 슬라이드 이미지 분석 (Phase 1: 차트 지원 추가)"""
        import requests

        # Phase 1: 차트 확인
        chart_info = self._extract_chart_info(slide)
        has_chart = chart_info["has_chart"]

        # 슬라이드를 PNG 이미지로 변환 (이미지가 제공되지 않은 경우만)
        if slide_img_base64 is None:
            try:
                slide_img_base64 = self._slide_to_base64_image(slide, slide_index)
            except Exception as e:
                raise RuntimeError(f"슬라이드 이미지 변환 실패: {e}")
        
        # Vision LLM API 호출 (멀티모달 지원)
        if llm_api_type == "openai":
            # OpenAI 공식 Vision API (base_url 무시)
            endpoint = "https://api.openai.com/v1/chat/completions"

            # 슬라이드 제목이 실제 제목인지 플레이스홀더인지 확인
            is_placeholder_title = slide_title.startswith("제목없음")

            # 차트 정보 추가 (Phase 1)
            chart_info_text = ""
            if has_chart:
                chart_types = ", ".join([c["type"] for c in chart_info["charts"]])
                chart_info_text = f"\n차트 정보: {chart_info['chart_count']}개 차트 ({chart_types})"

            # 제목 유무에 따라 다른 프롬프트 생성
            if is_placeholder_title:
                # 제목 없는 슬라이드용 프롬프트
                prompt_text = f"""[슬라이드 정보]
슬라이드: {slide_num}/{total_slides}
주의: 이 슬라이드에는 명시적인 제목이 없습니다{chart_info_text}

[중요 지시사항]
1. 슬라이드 번호 {slide_num}의 내용만 분석하세요
2. 다른 슬라이드와 혼동하지 마세요
3. 이미지에서 가장 큰 텍스트를 주제로 사용하세요

[분석 요구사항]
다음 정보를 구조화된 형식으로 추출하세요:

1. 주제: 슬라이드의 핵심 메시지 (1문장)
2. 데이터 유형: 표/차트/그래프 형태
   - 표: "N행 M열 표" 형식으로 명시
   - 차트: "막대/선/파이/영역 차트" 유형 명시
   - 차트가 있으면 트렌드, 이상치, 비교값도 분석
3. 주요 수치: 모든 숫자를 "항목명: 값 (단위)" 형식으로
4. 비교/추이: 증감률, 변화 패턴, 기간별 비교

[출력 형식]
```
주제: [핵심 메시지]
데이터 유형: [표/차트 상세]
주요 수치:
- [항목1]: [값1]
- [항목2]: [값2]
비교/추이: [...]
```

[주의사항]
- 다른 슬라이드 내용 혼동 금지
- 차트의 모든 데이터 포인트 추출
- 불명확한 수치는 [약]/[추정] 표시"""
            else:
                # 제목 있는 슬라이드용 프롬프트
                prompt_text = f"""[슬라이드 정보]
제목: "{slide_title}"
슬라이드: {slide_num}/{total_slides}{chart_info_text}

[중요 지시사항]
1. 슬라이드 맨 위의 큰 제목이 "{slide_title}"과 일치하는지 확인하세요
2. 다른 슬라이드의 내용을 분석하지 마세요
3. 위 슬라이드만 분석하세요

[분석 요구사항]
다음 정보를 구조화된 형식으로 추출하세요:

1. 주제: 슬라이드 제목 "{slide_title}" 기반 핵심 메시지 (1문장)
2. 데이터 유형: 표/차트/그래프 형태
   - 표: "N행 M열 표" 형식으로 명시
   - 차트: "막대/선/파이/영역 차트" 유형 명시
   - 차트가 있으면 트렌드, 이상치, 비교값도 분석
3. 주요 수치: 모든 숫자를 "항목명: 값 (단위)" 형식으로
4. 비교/추이: 증감률, 변화 패턴, 기간별 비교

[출력 형식]
```
주제: [슬라이드 제목 기반 메시지]
데이터 유형: [표/차트 상세]
주요 수치:
- [항목1]: [값1]
- [항목2]: [값2]
비교/추이: [...]
```

[예시]
주제: 2024 매출 성장 분석
데이터 유형: 4개 분기 비교 표 (4행 3열)
주요 수치:
- Q1 온라인: 150억원
- Q2 온라인: 180억원
- Q3 온라인: 190억원
- Q4 온라인: 195억원
비교/추이: Q4/Q1 +30% 성장

[주의사항]
- 슬라이드 제목과 일치하지 않으면 오류 보고
- 다른 슬라이드 내용 혼동 금지
- 차트의 모든 데이터 포인트 추출
- 불명확한 수치는 [약]/[추정] 표시"""

            # Phase 1: 차트가 있으면 detail "high" 사용
            detail_level = "high" if has_chart else "low"

            payload = {
                "model": llm_model,
                "messages": [{
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": prompt_text
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{slide_img_base64}",
                                "detail": detail_level  # 차트 있으면 high, 없으면 low
                            }
                        }
                    ]
                }],
                "max_tokens": 8000,  # 500→1000→8000 (Llama4-scout: 전체 슬라이드 상세 분석)
                "temperature": 0.1  # 일관성 최우선
            }
            
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {llm_api_key}" if llm_api_key else ""
            }
            
        elif llm_api_type in ["ollama", "request"]:
            # 비전 임베딩 설정 로드
            from config import ConfigManager
            cfg = ConfigManager().get_all()
            vision_enabled = cfg.get("vision_enabled", True)
            vision_mode = cfg.get("vision_mode", "auto")  # auto | ollama | openai-compatible
            
            # 헤더 확인을 위한 휴리스틱 함수
            def _is_openai_like(url: str, headers_dict: dict):
                return "/v1" in url or (headers_dict and "Authorization" in headers_dict and "Bearer" in headers_dict.get("Authorization", ""))
            
            # 비전 모드 결정
            def use_openai_style():
                if vision_mode == "openai-compatible":
                    return True
                if vision_mode == "ollama":
                    return False
                # auto: 휴리스틱으로 판단
                return _is_openai_like(llm_base_url, {})

            # 슬라이드 제목이 실제 제목인지 플레이스홀더인지 확인
            is_placeholder_title = slide_title.startswith("제목없음")

            # 제목 유무에 따라 다른 프롬프트 생성
            if is_placeholder_title:
                # 제목 없는 슬라이드용 프롬프트
                prompt_text = f"""[INST] [슬라이드 정보]
슬라이드: {slide_num}/{total_slides}
주의: 이 슬라이드에는 명시적인 제목이 없습니다

[중요 지시사항]
1. 슬라이드 번호 {slide_num}의 내용만 분석하세요
2. 다른 슬라이드와 혼동하지 마세요
3. 이미지에서 가장 큰 텍스트를 주제로 사용하세요

[분석 요구사항]
다음 정보를 구조화된 형식으로 추출하세요:

1. 주제: 슬라이드의 핵심 메시지 (1문장)
2. 데이터 유형: 표/그래프 형태
   - 표: "N행 M열 표" 형식으로 명시
   - 그래프: "막대/선/파이 그래프" 유형 명시
3. 주요 수치: 모든 숫자를 "항목명: 값 (단위)" 형식으로
4. 비교/추이: 증감률, 변화 패턴, 기간별 비교

[출력 형식]
```
주제: [핵심 메시지]
데이터 유형: [표/그래프 상세]
주요 수치:
- [항목1]: [값1]
- [항목2]: [값2]
비교/추이: [...]
```

[주의사항]
- 다른 슬라이드 내용 혼동 금지
- 불명확한 수치는 [약]/[추정] 표시 [/INST]"""
            else:
                # 제목 있는 슬라이드용 프롬프트 (기존 로직)
                prompt_text = f"""[INST] [슬라이드 정보]
제목: "{slide_title}"
슬라이드: {slide_num}/{total_slides}

[중요 지시사항]
1. 슬라이드 맨 위의 큰 제목이 "{slide_title}"과 일치하는지 확인하세요
2. 다른 슬라이드의 내용을 분석하지 마세요
3. 위 슬라이드만 분석하세요

[분석 요구사항]
다음 정보를 구조화된 형식으로 추출하세요:

1. 주제: 슬라이드 제목 "{slide_title}" 기반 핵심 메시지 (1문장)
2. 데이터 유형: 표/그래프 형태
   - 표: "N행 M열 표" 형식으로 명시
   - 그래프: "막대/선/파이 그래프" 유형 명시
3. 주요 수치: 모든 숫자를 "항목명: 값 (단위)" 형식으로
4. 비교/추이: 증감률, 변화 패턴, 기간별 비교

[출력 형식]
```
주제: [슬라이드 제목 기반 메시지]
데이터 유형: [표/그래프 상세]
주요 수치:
- [항목1]: [값1]
- [항목2]: [값2]
비교/추이: [...]
```

[예시]
주제: 2024 매출 성장 분석
데이터 유형: 4개 분기 비교 표 (4행 3열)
주요 수치:
- Q1 온라인: 150억원
- Q2 온라인: 180억원
- Q3 온라인: 190억원
- Q4 온라인: 195억원
비교/추이: Q4/Q1 +30% 성장

[주의사항]
- 슬라이드 제목과 일치하지 않으면 오류 보고
- 다른 슬라이드 내용 혼동 금지
- 불명확한 수치는 [약]/[추정] 표시 [/INST]"""
            
            # 비전이 활성화되고 이미지가 있는 경우
            if vision_enabled and slide_img_base64:
                # OpenAI 호환 방식 사용 여부 확인
                if use_openai_style():
                    # OpenAI 호환 멀티모달 API
                    endpoint = f"{llm_base_url.rstrip('/')}/v1/chat/completions"
                    payload = {
                        "model": llm_model,
                        "stream": False,
                        "messages": [{
                            "role": "user",
                            "content": [
                                {
                                    "type": "text",
                                    "text": prompt_text.replace("[INST]", "").replace("[/INST]", "").strip()
                                },
                                {
                                    "type": "image_url",
                                    "image_url": {
                                        "url": f"data:image/png;base64,{slide_img_base64}",
                                        "detail": "low"
                                    }
                                }
                            ]
                        }],
                        "max_tokens": 8000,  # 1000 → 8000 (Llama4-scout: 숨겨진 슬라이드 상세 분석)
                        "temperature": 0.1
                    }
                    headers = {
                        "Content-Type": "application/json",
                        "Authorization": f"Bearer {llm_api_key}" if llm_api_key else ""
                    }
                else:
                    # Ollama 네이티브 멀티모달 API
                    endpoint = f"{llm_base_url.rstrip('/')}/api/chat"
                    payload = {
                        "model": llm_model,
                        "stream": False,
                        "messages": [{
                            "role": "user",
                            "content": prompt_text,
                            "images": [slide_img_base64]  # 순수 base64 (data: 접두 없음)
                        }],
                        "options": {
                            "temperature": 0.1,
                            "num_predict": 1000
                        }
                    }
                    headers = {"Content-Type": "application/json"}
            else:
                # 비전 비활성화 또는 이미지 없음 - 기존 텍스트 전용 방식
                endpoint = f"{llm_base_url.rstrip('/')}/api/generate"
                payload = {
                    "model": llm_model,
                    "prompt": prompt_text,
                    "stream": False,
                    "options": {
                        "temperature": 0.1,
                        "num_predict": 1000
                    }
                }
                headers = {"Content-Type": "application/json"}
        
        else:
            raise ValueError(f"지원하지 않는 Vision API 타입: {llm_api_type}")
        
        try:
            print(f"[Vision] 슬라이드 Vision 분석 요청: {llm_api_type} - {llm_model}")
            response = requests.post(endpoint, json=payload, headers=headers, timeout=60)
            
            if response.status_code == 200:
                result = response.json()
                
                if llm_api_type == "openai":
                    vision_text = result["choices"][0]["message"]["content"]
                elif llm_api_type in ["ollama", "request"]:
                    # 비전 모드 확인
                    from config import ConfigManager
                    cfg = ConfigManager().get_all()
                    vision_mode = cfg.get("vision_mode", "auto")
                    
                    def use_openai_style():
                        if vision_mode == "openai-compatible":
                            return True
                        if vision_mode == "ollama":
                            return False
                        # auto: 엔드포인트로 판단
                        return "/v1/chat/completions" in endpoint
                    
                    if use_openai_style() and "/v1/chat/completions" in endpoint:
                        # OpenAI 호환 응답 처리
                        vision_text = result["choices"][0]["message"]["content"]
                    else:
                        # Ollama 네이티브 응답 처리
                        if "message" in result and "content" in result["message"]:
                            vision_text = result["message"]["content"]
                        elif "response" in result:
                            vision_text = result["response"]
                        else:
                            vision_text = str(result)
                else:  # ollama, request (기존 방식)
                    vision_text = result.get("response", "")
                
                print(f"[Vision] [OK] 분석 완료: {len(vision_text)}자")
                return vision_text
            else:
                raise RuntimeError(f"Vision API 오류 ({response.status_code}): {response.text}")
                
        except requests.exceptions.RequestException as e:
            raise RuntimeError(f"Vision API 네트워크 오류: {e}")
    
    def _slide_to_base64_image(self, slide, slide_index: int) -> str:
        """슬라이드를 Base64 PNG 이미지로 변환"""
        import sys
        
        try:
            # 방법 1: Windows COM 사용 (PowerPoint 활성화 필요)
            if sys.platform == "win32":
                try:
                    return self._slide_to_image_via_com(slide_index)
                except Exception as e:
                    print(f"[Vision] COM 방식 실패: {e}, Pillow 방식 시도")
            
            # 방법 2: Pillow로 간단한 렌더링 (폴백)
            from PIL import Image, ImageDraw, ImageFont
            
            # 슬라이드 크기 가져오기 (presentation에서)
            if hasattr(self, 'presentation'):
                # EMU 단위를 픽셀로 변환 (9144000 EMU = 10 inches = 960 pixels @ 96 DPI)
                slide_width_emu = self.presentation.slide_width
                slide_height_emu = self.presentation.slide_height
                slide_width = int(slide_width_emu / 9144000 * 960)
                slide_height = int(slide_height_emu / 9144000 * 960)
            else:
                # 기본 크기
                slide_width = 960
                slide_height = 540
            
            # 빈 이미지 생성 (배경은 흰색)
            img = Image.new('RGB', (slide_width, slide_height), 'white')
            draw = ImageDraw.Draw(img)
            
            # 슬라이드의 텍스트를 간단히 렌더링
            try:
                # 제목 렌더링
                if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
                    title_text = slide.shapes.title.text
                    if title_text:
                        font = ImageFont.load_default()
                        draw.text((50, 50), title_text, fill='black', font=font)
                
                # 본문 텍스트 렌더링 (처음 500자만)
                full_text = self._extract_full_text_from_slide(slide)
                if full_text:
                    font = ImageFont.load_default()
                    # 긴 텍스트를 여러 줄로 분할
                    lines = full_text[:500].split('\n')[:20]  # 최대 20줄
                    y_offset = 150
                    for line in lines:
                        if line.strip():
                            draw.text((50, y_offset), line[:80], fill='black', font=font)
                            y_offset += 30
            except Exception as e:
                print(f"[Vision] 텍스트 렌더링 경고: {e}")
            
            buffer = io.BytesIO()
            img.save(buffer, format='PNG')
            img_bytes = buffer.getvalue()
            
            return base64.b64encode(img_bytes).decode('utf-8')
            
        except ImportError:
            # Pillow 미설치 시 빈 문자열 반환
            print("[Vision] Pillow 미설치: 슬라이드 이미지 변환 건너뜀")
            raise RuntimeError("Pillow 라이브러리가 필요합니다")
        except Exception as e:
            print(f"[Vision] 슬라이드 이미지 변환 실패: {e}")
            raise RuntimeError(f"슬라이드 이미지 변환 실패: {e}")
    
    def _render_all_slides_via_com(self, total_slides: int) -> Dict[int, dict]:
        """Windows COM을 사용하여 모든 슬라이드를 한 번에 이미지로 변환 (Windows 전용)

        Returns:
            Dict[int, dict]: {index: {"image": base64_str, "title": str, "slide_id": int}}
        """
        import win32com.client
        from PIL import Image
        import os
        import tempfile

        slide_images = {}
        
        try:
            if not hasattr(self, 'pptx_path'):
                raise RuntimeError("pptx_path가 설정되지 않았습니다")
            
            # PowerPoint 애플리케이션 시작 (한 번만)
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            # PowerPoint는 Visible=False를 지원하지 않으므로 표시 상태 유지
            
            try:
                # 프레젠테이션 열기 (절대 경로로 변환)
                abs_path = os.path.abspath(self.pptx_path)
                presentation = powerpoint.Presentations.Open(abs_path)
                
                print(f"[Vision] PowerPoint 열림: {total_slides}개 슬라이드")
                
                # 디버그 디렉토리 생성 (환경변수로 활성화)
                debug_vision = os.getenv("DEBUG_VISION_IMAGES", "false").lower() == "true"
                debug_dir = None
                if debug_vision:
                    debug_dir = os.path.join(os.path.dirname(abs_path), "debug_vision")
                    os.makedirs(debug_dir, exist_ok=True)
                    print(f"[Vision] [DEBUG] 이미지 저장 위치: {debug_dir}")

                # 모든 슬라이드를 한 번에 렌더링
                for slide_index in range(total_slides):
                    try:
                        # 슬라이드 선택 (1-based index)
                        slide = presentation.Slides[slide_index + 1]

                        # 슬라이드 제목 추출
                        slide_title = ""
                        try:
                            if slide.Shapes.HasTitle:
                                slide_title = slide.Shapes.Title.TextFrame.TextRange.Text.strip()
                        except:
                            pass

                        # 제목이 없으면 슬라이드 번호로 대체 (python-pptx와 동일한 형식)
                        if not slide_title:
                            slide_title = f"제목없음-슬라이드{slide_index + 1}"

                        # 슬라이드 ID 추출
                        slide_id = slide.SlideID

                        # 임시 파일 경로
                        temp_file = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
                        temp_path = temp_file.name
                        temp_file.close()

                        # 슬라이드를 PNG로 내보내기
                        slide.Export(temp_path, "PNG", 1920, 1080)  # 1920x1080 해상도

                        # 이미지 읽기
                        img = Image.open(temp_path)
                        buffer = io.BytesIO()
                        img.save(buffer, format='PNG')
                        img_bytes = buffer.getvalue()

                        # Base64 변환 및 이미지 + 메타데이터 저장
                        slide_images[slide_index] = {
                            "image": base64.b64encode(img_bytes).decode('utf-8'),
                            "title": slide_title,
                            "slide_id": slide_id,
                            "com_index": slide_index
                        }

                        # 디버그: 이미지 저장 (환경변수로 활성화 시)
                        if debug_vision and debug_dir:
                            title_suffix = f"_{slide_title[:20]}" if slide_title else "_notitle"
                            debug_filename = f"slide_{slide_index+1}_com_index_{slide_index}{title_suffix}.png"
                            debug_path = os.path.join(debug_dir, debug_filename)
                            img.save(debug_path)
                            print(f"  [DEBUG] 저장: {debug_filename} (제목: {slide_title if slide_title else '없음'})")

                        # 임시 파일 삭제
                        os.unlink(temp_path)

                        print(f"  [OK] 슬라이드 {slide_index + 1} 렌더링 완료 (제목: {slide_title if slide_title else '없음'})")
                    except Exception as e:
                        print(f"  [WARN] 슬라이드 {slide_index + 1} 렌더링 실패: {e}")
                
                # 프레젠테이션 닫기
                presentation.Close()
                
            finally:
                powerpoint.Quit()
                print("[Vision] PowerPoint 종료됨")
                
        except ImportError:
            raise RuntimeError("pywin32가 설치되지 않았습니다 (Windows 전용)")
        except Exception as e:
            raise RuntimeError(f"COM 이미지 변환 실패: {e}")
        
        return slide_images
    
    def _slide_to_image_via_com(self, slide_index: int) -> str:
        """Windows COM을 사용하여 슬라이드를 이미지로 변환 (Windows 전용) - 레거시 메서드"""
        import win32com.client
        from PIL import Image
        
        try:
            if not hasattr(self, 'pptx_path'):
                raise RuntimeError("pptx_path가 설정되지 않았습니다")
            
            # PowerPoint 애플리케이션 시작
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            # PowerPoint는 Visible=False를 지원하지 않으므로 표시 상태 유지
            
            try:
                # 프레젠테이션 열기 (절대 경로로 변환)
                import os
                abs_path = os.path.abspath(self.pptx_path)
                presentation = powerpoint.Presentations.Open(abs_path)
                
                # 슬라이드 선택 (1-based index)
                slide = presentation.Slides[slide_index + 1]
                
                # 임시 파일 경로
                import tempfile
                temp_file = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
                temp_path = temp_file.name
                temp_file.close()
                
                # 슬라이드를 PNG로 내보내기
                slide.Export(temp_path, "PNG", 1920, 1080)  # 1920x1080 해상도
                
                # 이미지 읽기
                img = Image.open(temp_path)
                buffer = io.BytesIO()
                img.save(buffer, format='PNG')
                img_bytes = buffer.getvalue()
                
                # 프레젠테이션 닫기
                presentation.Close()
                
                # 임시 파일 삭제
                import os
                os.unlink(temp_path)
                
                return base64.b64encode(img_bytes).decode('utf-8')
                
            finally:
                powerpoint.Quit()
                
        except ImportError:
            raise RuntimeError("pywin32가 설치되지 않았습니다 (Windows 전용)")
        except Exception as e:
            raise RuntimeError(f"COM 이미지 변환 실패: {e}")
