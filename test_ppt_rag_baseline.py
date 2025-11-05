"""
PPT RAG 사전 테스트 (비전 청킹 제외)
Phase 1 & 2 구현 전 베이스라인 성능 측정
"""
import sys
import os
import time
import json
from pathlib import Path
from typing import List, Dict, Any
from pptx import Presentation

sys.path.insert(0, os.path.dirname(__file__))

from utils.pptx_chunking_engine import PPTXChunkingEngine


class PPTRAGBaselineTester:
    """PPT RAG 베이스라인 테스트 클래스"""

    def __init__(self, test_pptx_dir: str = "data/test_pptx"):
        self.test_pptx_dir = Path(test_pptx_dir)
        self.results = []

    def select_test_files(self, count: int = 3) -> List[Path]:
        """테스트용 PPT 파일 선택 (표/그래프가 많은 복잡한 파일 우선)"""
        priority_files = [
            "complex_03_data_analysis_report.pptx",
            "advanced_01_financial_report.pptx",
            "complex_04_project_plan.pptx",
        ]

        selected = []
        for filename in priority_files[:count]:
            filepath = self.test_pptx_dir / filename
            if filepath.exists():
                selected.append(filepath)

        # 우선순위 파일이 부족하면 다른 파일 추가
        if len(selected) < count:
            all_pptx = list(self.test_pptx_dir.glob("*.pptx"))
            for pptx_file in all_pptx:
                if pptx_file not in selected:
                    selected.append(pptx_file)
                    if len(selected) >= count:
                        break

        return selected

    def analyze_ppt_structure(self, pptx_path: Path) -> Dict[str, Any]:
        """PPT 파일 구조 분석 (슬라이드, 표, 텍스트 등)"""
        prs = Presentation(str(pptx_path))

        slide_count = len(prs.slides)
        table_count = 0
        text_shape_count = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'table'):
                    try:
                        if shape.has_table:
                            table_count += 1
                    except:
                        pass
                if hasattr(shape, 'text_frame'):
                    try:
                        if shape.has_text_frame and shape.text.strip():
                            text_shape_count += 1
                    except:
                        pass

        return {
            "slide_count": slide_count,
            "table_count": table_count,
            "text_shape_count": text_shape_count
        }

    def test_single_pptx(self, pptx_path: Path) -> Dict[str, Any]:
        """단일 PPT 파일 테스트 (비전 청킹 제외)"""
        print(f"\n{'='*80}")
        print(f"[TEST] {pptx_path.name}")
        print(f"{'='*80}")

        # 파일 구조 분석
        structure = self.analyze_ppt_structure(pptx_path)
        print(f"\n[구조 분석]")
        print(f"  - 슬라이드: {structure['slide_count']}개")
        print(f"  - 표: {structure['table_count']}개")
        print(f"  - 텍스트 Shape: {structure['text_shape_count']}개")

        # 청킹 엔진 초기화
        config = {
            "max_size": 300,
            "overlap_size": 50,
            "enable_small_to_large": True
        }
        chunker = PPTXChunkingEngine(config)

        # 시작 시간 측정
        start_time = time.time()

        try:
            # PPT 처리 (비전 청킹 비활성화)
            chunks = chunker.process_pptx_document(
                pptx_path=str(pptx_path),
                enable_vision=False  # 비전 청킹 제외
            )

            # 종료 시간 측정
            elapsed_time = time.time() - start_time

            # 청크 분석
            text_chunks = len(chunks)
            table_chunks = sum(1 for chunk in chunks
                             if chunk.chunk_type == 'table')
            bullet_chunks = sum(1 for chunk in chunks
                              if chunk.chunk_type == 'bullet_group')
            slide_summary_chunks = sum(1 for chunk in chunks
                                     if chunk.chunk_type == 'slide_summary')

            # 평균 청크 크기
            avg_chunk_size = sum(len(chunk.content) for chunk in chunks) / len(chunks) if chunks else 0

            # 결과 저장
            result = {
                "file_name": pptx_path.name,
                "structure": structure,
                "chunk_count": text_chunks,
                "chunk_types": {
                    "table": table_chunks,
                    "bullet": bullet_chunks,
                    "slide_summary": slide_summary_chunks,
                    "other": text_chunks - table_chunks - bullet_chunks - slide_summary_chunks
                },
                "avg_chunk_size": avg_chunk_size,
                "total_time": elapsed_time,
                "time_per_slide": elapsed_time / structure['slide_count'] if structure['slide_count'] > 0 else 0,
                "success": True,
                "error": None
            }

            print(f"\n[OK] 처리 완료:")
            print(f"  - 총 청크 수: {text_chunks}개")
            print(f"  - 표 청크: {table_chunks}개")
            print(f"  - 불릿 청크: {bullet_chunks}개")
            print(f"  - 슬라이드 요약: {slide_summary_chunks}개")
            print(f"  - 평균 청크 크기: {avg_chunk_size:.0f}자")
            print(f"  - 총 소요 시간: {elapsed_time:.2f}초")
            print(f"  - 슬라이드당 시간: {result['time_per_slide']:.2f}초")

            # 샘플 청크 출력 (표 청크 우선)
            print(f"\n[샘플 청크]")
            sample_count = 0
            for chunk in chunks:
                if sample_count >= 2:
                    break
                if chunk.chunk_type == 'table':
                    print(f"\n  [표 청크 예시]")
                    print(f"  {chunk.content[:200]}...")
                    sample_count += 1

            return result

        except Exception as e:
            elapsed_time = time.time() - start_time

            result = {
                "file_name": pptx_path.name,
                "structure": structure,
                "chunk_count": 0,
                "chunk_types": {},
                "avg_chunk_size": 0,
                "total_time": elapsed_time,
                "time_per_slide": 0,
                "success": False,
                "error": str(e)
            }

            print(f"\n[NG] 처리 실패: {e}")
            import traceback
            traceback.print_exc()

            return result

    def run_baseline_suite(self) -> List[Dict[str, Any]]:
        """베이스라인 테스트 스위트 실행"""
        print(f"\n{'='*80}")
        print(f"PPT RAG 베이스라인 테스트 (비전 청킹 제외)")
        print(f"{'='*80}")

        # 테스트 파일 선택
        test_files = self.select_test_files(count=3)

        if not test_files:
            print("[NG] 테스트 파일을 찾을 수 없습니다.")
            return []

        print(f"\n테스트 파일 {len(test_files)}개:")
        for f in test_files:
            print(f"  - {f.name}")

        # 각 파일 테스트
        results = []
        for pptx_path in test_files:
            result = self.test_single_pptx(pptx_path)
            results.append(result)

        return results

    def print_summary(self, results: List[Dict[str, Any]]):
        """테스트 결과 요약"""
        print(f"\n{'='*80}")
        print(f"베이스라인 테스트 결과 요약")
        print(f"{'='*80}")

        if not results:
            print("테스트 결과가 없습니다.")
            return

        # 성공/실패 개수
        success_count = sum(1 for r in results if r['success'])
        total_count = len(results)

        print(f"\n전체: {total_count}개 파일")
        print(f"성공: {success_count}개")
        print(f"실패: {total_count - success_count}개")

        # 성공한 테스트만 집계
        success_results = [r for r in results if r['success']]

        if not success_results:
            print("\n성공한 테스트가 없습니다.")
            return

        # 평균 계산
        total_slides = sum(r['structure']['slide_count'] for r in success_results)
        total_tables = sum(r['structure']['table_count'] for r in success_results)
        total_chunks = sum(r['chunk_count'] for r in success_results)
        total_table_chunks = sum(r['chunk_types'].get('table', 0) for r in success_results)
        total_time = sum(r['total_time'] for r in success_results)
        avg_time_per_slide = total_time / total_slides if total_slides > 0 else 0
        avg_chunk_size = sum(r['avg_chunk_size'] for r in success_results) / len(success_results)

        print(f"\n[STAT] 집계:")
        print(f"  - 총 슬라이드: {total_slides}개")
        print(f"  - 총 표: {total_tables}개")
        print(f"  - 총 청크: {total_chunks}개")
        print(f"  - 표 청크: {total_table_chunks}개 (표당 {total_table_chunks/total_tables:.1f}개)" if total_tables > 0 else "  - 표 청크: 0개")
        print(f"  - 평균 청크 크기: {avg_chunk_size:.0f}자")
        print(f"  - 총 소요 시간: {total_time:.2f}초")
        print(f"  - 평균 슬라이드당 시간: {avg_time_per_slide:.2f}초")

        # 개별 결과
        print(f"\n[FILE] 개별 결과:")
        for r in success_results:
            print(f"\n  [{r['file_name']}]")
            print(f"    - 슬라이드: {r['structure']['slide_count']}개")
            print(f"    - 표: {r['structure']['table_count']}개")
            print(f"    - 청크: {r['chunk_count']}개 (표: {r['chunk_types'].get('table', 0)}개)")
            print(f"    - 시간: {r['total_time']:.2f}초 ({r['time_per_slide']:.2f}초/슬라이드)")

        return {
            "total_slides": total_slides,
            "total_tables": total_tables,
            "total_chunks": total_chunks,
            "total_table_chunks": total_table_chunks,
            "avg_chunk_size": avg_chunk_size,
            "total_time": total_time,
            "avg_time_per_slide": avg_time_per_slide,
            "success_rate": success_count / total_count if total_count > 0 else 0
        }


def main():
    """메인 실행 함수"""
    print("\n" + "="*80)
    print("PPT RAG 베이스라인 테스트 시작")
    print("Phase 1 & 2 구현 전 현재 성능 측정")
    print("="*80)

    tester = PPTRAGBaselineTester()

    # 테스트 실행
    results = tester.run_baseline_suite()

    # 결과 요약
    summary = tester.print_summary(results)

    # 결과 저장 (JSON)
    output_file = "test_results/ppt_rag_baseline.json"
    os.makedirs("test_results", exist_ok=True)

    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump({
            "test_name": "PPT RAG Baseline (No Vision)",
            "timestamp": time.strftime("%Y-%m-%d %H:%M:%S"),
            "results": results,
            "summary": summary
        }, f, indent=2, ensure_ascii=False)

    print(f"\n결과 저장: {output_file}")

    return results, summary


if __name__ == "__main__":
    main()
    sys.exit(0)
