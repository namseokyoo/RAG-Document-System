"""
Vision 청킹 통합 테스트 - 여러 PPT 파일 검증

Phase 1-B 구현 후 다양한 PPT 파일로 Vision 청킹 성능 검증
"""

import sys
import os
from pathlib import Path


def test_vision_integration():
    """여러 PPT 파일로 Vision 청킹 통합 테스트"""

    print("=" * 80)
    print("Vision 청킹 통합 테스트 - 다중 파일 검증")
    print("=" * 80)
    print()

    # 테스트 파일 목록
    test_files = [
        "data/test_pptx/advanced_01_financial_report.pptx",      # 재무 보고서 (이미 테스트)
        "data/test_pptx/advanced_02_product_plan.pptx",          # 제품 계획
        "data/test_pptx/complex_03_data_analysis_report.pptx",   # 데이터 분석
        "data/test_pptx/complex_05_comprehensive_report.pptx",   # 종합 리포트
    ]

    # 환경 설정
    from utils.pptx_chunking_engine import PPTXChunkingEngine
    from config import ConfigManager

    config_mgr = ConfigManager()
    config = config_mgr.get_all()

    llm_api_type = config.get("llm_api_type", "openai")
    llm_model = config.get("llm_model", "gpt-4o-mini")
    llm_api_key = config.get("llm_api_key", os.getenv("OPENAI_API_KEY"))

    print(f"설정:")
    print(f"  API 타입: {llm_api_type}")
    print(f"  모델: {llm_model}")
    print(f"  Vision 활성화: {config.get('enable_vision_chunking', True)}")
    print()

    # 청킹 엔진 초기화
    engine = PPTXChunkingEngine(config)

    # 결과 저장
    results = []

    # 각 파일 테스트
    for idx, file_path in enumerate(test_files, 1):
        test_file = Path(file_path)

        if not test_file.exists():
            print(f"[{idx}/{len(test_files)}] SKIP: {test_file.name} (파일 없음)")
            print()
            continue

        print("=" * 80)
        print(f"[{idx}/{len(test_files)}] 테스트: {test_file.name}")
        print("=" * 80)
        print()

        try:
            # PPT 파일 기본 정보
            from pptx import Presentation
            prs = Presentation(str(test_file))

            slide_count = len(prs.slides)
            print(f"총 슬라이드 수: {slide_count}")

            # 각 슬라이드 구조 확인
            print()
            print("슬라이드 구조:")
            print("-" * 80)

            title_count = 0
            titleless_count = 0
            table_count = 0

            for i, slide in enumerate(prs.slides, 1):
                # 제목 추출
                title = ""
                try:
                    if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
                        title = slide.shapes.title.text.strip()
                except:
                    pass

                if not title:
                    title = f"제목없음-슬라이드{i}"
                    titleless_count += 1
                else:
                    title_count += 1

                # 표 개수
                slide_tables = 0
                for shape in slide.shapes:
                    if shape.has_table:
                        slide_tables += 1
                        table_count += 1

                print(f"  슬라이드 {i}: {title}")
                if slide_tables > 0:
                    print(f"    표: {slide_tables}개")

            print()
            print(f"요약: {slide_count}개 슬라이드 (제목 있음: {title_count}, 없음: {titleless_count}, 표: {table_count}개)")
            print()

            # Vision 청킹 실행
            print("Vision 청킹 실행 중...")
            print("-" * 80)

            chunks = engine.process_pptx_document(
                pptx_path=str(test_file.absolute()),
                enable_vision=True,
                llm_api_type=llm_api_type,
                llm_base_url=config.get('llm_base_url', ''),
                llm_model=llm_model,
                llm_api_key=llm_api_key
            )

            if not chunks:
                print("[ERROR] 청킹 실패 - 청크 0개")
                results.append({
                    "file": test_file.name,
                    "status": "FAIL",
                    "slides": slide_count,
                    "chunks": 0,
                    "vision_used": 0,
                    "error": "청킹 실패"
                })
                print()
                continue

            # 결과 분석
            print(f"총 {len(chunks)}개 청크 생성")
            print()

            # Vision 사용 여부 확인
            vision_chunks = 0
            slide_summary_chunks = []

            for chunk in chunks:
                # slide_summary 타입 청크만 확인
                if chunk.chunk_type == "slide_summary":
                    slide_summary_chunks.append(chunk)
                    # Vision으로 생성된 청크인지 확인 (content에 "주제:" 패턴이 있으면 Vision)
                    if "주제:" in chunk.content or "데이터 유형:" in chunk.content:
                        vision_chunks += 1

            print(f"슬라이드 요약 청크: {len(slide_summary_chunks)}개")
            print(f"Vision 사용: {vision_chunks}개")
            print(f"Vision 사용률: {vision_chunks}/{len(slide_summary_chunks)} = {vision_chunks/len(slide_summary_chunks)*100:.1f}%")
            print()

            # 샘플 출력 (첫 3개 슬라이드)
            print("샘플 결과 (최대 3개):")
            print("-" * 80)

            for i, chunk in enumerate(slide_summary_chunks[:3], 1):
                slide_num = chunk.metadata.slide_number
                slide_title = chunk.metadata.slide_title

                is_vision = "주제:" in chunk.content or "데이터 유형:" in chunk.content
                method = "[Vision]" if is_vision else "[Text]"

                print(f"\n슬라이드 {slide_num}: {slide_title} {method}")
                print(f"  내용 (앞 200자):")
                print(f"  {chunk.content[:200]}...")

            print()
            print("-" * 80)
            print()

            # 결과 저장
            results.append({
                "file": test_file.name,
                "status": "SUCCESS",
                "slides": slide_count,
                "titleless": titleless_count,
                "tables": table_count,
                "chunks": len(chunks),
                "slide_summaries": len(slide_summary_chunks),
                "vision_used": vision_chunks,
                "vision_rate": f"{vision_chunks/len(slide_summary_chunks)*100:.1f}%"
            })

        except Exception as e:
            print(f"[ERROR] 테스트 실패: {e}")
            import traceback
            traceback.print_exc()

            results.append({
                "file": test_file.name,
                "status": "ERROR",
                "error": str(e)
            })
            print()

    # 최종 요약
    print()
    print("=" * 80)
    print("전체 테스트 결과 요약")
    print("=" * 80)
    print()

    success_count = sum(1 for r in results if r["status"] == "SUCCESS")
    print(f"테스트 파일: {len(results)}개")
    print(f"성공: {success_count}개")
    print(f"실패: {len(results) - success_count}개")
    print()

    # 테이블 형식 출력
    print(f"{'파일명':<50} {'슬라이드':<8} {'제목없음':<8} {'표':<6} {'Vision사용':<12} {'상태'}")
    print("-" * 100)

    for r in results:
        if r["status"] == "SUCCESS":
            print(f"{r['file']:<50} {r['slides']:<8} {r['titleless']:<8} {r['tables']:<6} "
                  f"{r['vision_used']}/{r['slide_summaries']:<5} ({r['vision_rate']:<5}) SUCCESS")
        else:
            print(f"{r['file']:<50} {'N/A':<8} {'N/A':<8} {'N/A':<6} {'N/A':<12} {r['status']}")

    print()
    print("=" * 80)
    print("테스트 완료")
    print("=" * 80)


if __name__ == "__main__":
    try:
        test_vision_integration()
    except Exception as e:
        print(f"[ERROR] 테스트 실패: {e}")
        import traceback
        traceback.print_exc()
