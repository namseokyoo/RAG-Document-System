"""
차트 테스트용 PPT 생성 스크립트
"""
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

def create_chart_test_ppt():
    """차트가 포함된 테스트 PPT 생성"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 슬라이드 1: 막대 차트
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout

    # 제목 설정 (title placeholder 사용)
    title1 = slide1.shapes.title
    title1.text = "2024년 분기별 매출 현황"

    # 막대 차트 데이터
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
    chart_data.add_series('온라인', (150, 180, 190, 195))
    chart_data.add_series('오프라인', (120, 115, 110, 105))

    # 차트 추가
    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
    chart = slide1.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = "분기별 매출 비교 (단위: 억원)"

    print("  슬라이드 1: 막대 차트 추가 완료")

    # 슬라이드 2: 선 차트
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout

    # 제목 설정
    title2 = slide2.shapes.title
    title2.text = "월별 성장률 추이"

    # 선 차트 데이터
    chart_data2 = CategoryChartData()
    chart_data2.categories = ['1월', '2월', '3월', '4월', '5월', '6월']
    chart_data2.add_series('성장률(%)', (5.2, 7.5, 6.8, 9.1, 8.3, 10.5))

    # 차트 추가
    chart2 = slide2.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data2
    ).chart

    chart2.has_title = True
    chart2.chart_title.text_frame.text = "월별 성장률 변화"

    print("  슬라이드 2: 선 차트 추가 완료")

    # 슬라이드 3: 파이 차트
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout

    # 제목 설정
    title3 = slide3.shapes.title
    title3.text = "제품별 시장 점유율"

    # 파이 차트 데이터
    chart_data3 = CategoryChartData()
    chart_data3.categories = ['제품 A', '제품 B', '제품 C', '기타']
    chart_data3.add_series('점유율', (35, 28, 22, 15))

    # 차트 추가
    chart3 = slide3.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data3
    ).chart

    chart3.has_title = True
    chart3.chart_title.text_frame.text = "시장 점유율 분포 (%)"

    print("  슬라이드 3: 파이 차트 추가 완료")

    # 슬라이드 4: 차트 + 표 혼합
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout

    # 제목 설정
    title4 = slide4.shapes.title
    title4.text = "종합 분석 대시보드"

    # 작은 차트 추가
    chart_data4 = CategoryChartData()
    chart_data4.categories = ['A', 'B', 'C']
    chart_data4.add_series('값', (100, 150, 120))

    chart4 = slide4.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(2), Inches(4), Inches(3), chart_data4
    ).chart

    # 표 추가
    rows, cols = 4, 3
    table_shape = slide4.shapes.add_table(rows, cols, Inches(5.5), Inches(2), Inches(4), Inches(3))
    table = table_shape.table

    # 표 헤더
    table.cell(0, 0).text = "지역"
    table.cell(0, 1).text = "매출"
    table.cell(0, 2).text = "성장률"

    # 표 데이터
    data = [
        ("서울", "250억", "15%"),
        ("부산", "180억", "12%"),
        ("대구", "120억", "8%")
    ]

    for i, (region, sales, growth) in enumerate(data, 1):
        table.cell(i, 0).text = region
        table.cell(i, 1).text = sales
        table.cell(i, 2).text = growth

    print("  슬라이드 4: 차트 + 표 혼합 추가 완료")

    # 저장
    output_path = "data/test_pptx/chart_test.pptx"
    prs.save(output_path)
    print(f"\n[SUCCESS] 차트 테스트 PPT 생성 완료: {output_path}")
    print(f"   총 {len(prs.slides)}개 슬라이드")
    print(f"   - 슬라이드 1: 막대 차트 (분기별 매출)")
    print(f"   - 슬라이드 2: 선 차트 (월별 성장률)")
    print(f"   - 슬라이드 3: 파이 차트 (시장 점유율)")
    print(f"   - 슬라이드 4: 차트 + 표 혼합")

if __name__ == "__main__":
    print("=" * 80)
    print("차트 테스트 PPT 생성")
    print("=" * 80)
    print()
    create_chart_test_ppt()
    print()
    print("=" * 80)
