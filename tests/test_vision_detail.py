"""
Vision 청킹 상세 분석 - 실제 내용 확인
실제 PPT 내용과 Vision 분석 결과를 상세히 비교
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from utils.pptx_chunking_engine import PPTXChunkingEngine
from config import ConfigManager
from pptx import Presentation


def analyze_ppt_detailed():
    """PPT 파일의 실제 내용과 Vision 분석 결과 상세 비교"""

    test_file = Path("data/test_pptx/advanced_01_financial_report.pptx")

    print("=" * 80)
    print(f"상세 분석: {test_file.name}")
    print("=" * 80)

    # 1. 실제 PPT 내용 추출
    print("\n[1] 실제 PPT 내용 (Ground Truth)")
    print("-" * 80)

    prs = Presentation(str(test_file))

    for slide_idx, slide in enumerate(prs.slides, 1):
        print(f"\n슬라이드 {slide_idx}")
        print("=" * 80)

        # 제목
        if slide.shapes.title:
            print(f"제목: {slide.shapes.title.text}")

        # 모든 텍스트
        print("\n[텍스트 내용]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                print(f"  - {shape.text[:100]}")

        # 표 확인
        has_table = False
        for shape in slide.shapes:
            if shape.has_table:
                has_table = True
                table = shape.table
                print(f"\n[표 발견] {len(table.rows)}행 x {len(table.columns)}열")

                # 표 내용 출력
                for row_idx, row in enumerate(table.rows):
                    row_data = [cell.text for cell in row.cells]
                    print(f"  행 {row_idx}: {row_data}")

        if not has_table:
            print("\n[표 없음]")

        # 차트 확인
        has_chart = False
        for shape in slide.shapes:
            if shape.has_chart:
                has_chart = True
                print(f"\n[차트 발견] 타입: {shape.chart.chart_type}")

        if not has_chart:
            print("\n[차트 없음]")

    # 2. Vision 분석 결과
    print("\n\n" + "=" * 80)
    print("[2] Vision 분석 결과")
    print("=" * 80)

    config_mgr = ConfigManager()
    config = config_mgr.get_all()

    engine = PPTXChunkingEngine(config)

    chunks = engine.process_pptx_document(
        pptx_path=str(test_file),
        enable_vision=True,
        llm_api_type=config['llm_api_type'],
        llm_base_url=config.get('llm_base_url', ''),
        llm_model=config['llm_model'],
        llm_api_key=config.get('llm_api_key', '')
    )

    # 슬라이드별로 그룹핑
    slides_content = {}
    for chunk in chunks:
        slide_num = chunk.metadata.slide_number
        chunk_text = chunk.content if hasattr(chunk, 'content') else chunk.text if hasattr(chunk, 'text') else str(chunk)

        if slide_num not in slides_content:
            slides_content[slide_num] = {
                "slide_summary": None,
                "other_chunks": []
            }

        if chunk.chunk_type == "slide_summary":
            slides_content[slide_num]["slide_summary"] = chunk_text
        else:
            slides_content[slide_num]["other_chunks"].append({
                "type": chunk.chunk_type,
                "text": chunk_text
            })

    # Vision 분석 결과 출력
    for slide_num in sorted(slides_content.keys()):
        print(f"\n슬라이드 {slide_num}")
        print("=" * 80)

        content = slides_content[slide_num]

        if content["slide_summary"]:
            print("\n[슬라이드 요약 (Vision 포함)]")
            print(content["slide_summary"])

        print(f"\n[기타 청크] {len(content['other_chunks'])}개")
        for i, chunk in enumerate(content["other_chunks"][:3]):  # 처음 3개만
            print(f"\n  청크 {i+1} ({chunk['type']}):")
            print(f"  {chunk['text'][:200]}...")

    # 3. 비교 분석
    print("\n\n" + "=" * 80)
    print("[3] 비교 분석")
    print("=" * 80)

    # 슬라이드 3 집중 분석 (표가 있는 슬라이드)
    print("\n슬라이드 3 집중 분석 (표 포함)")
    print("-" * 80)

    slide_3 = prs.slides[2]  # 0-indexed

    print("\n[실제 표 내용]")
    for shape in slide_3.shapes:
        if shape.has_table:
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                row_data = [cell.text for cell in row.cells]
                print(f"  {' | '.join(row_data)}")

    print("\n[Vision 분석 결과]")
    if 3 in slides_content and slides_content[3]["slide_summary"]:
        vision_text = slides_content[3]["slide_summary"]

        # [Vision Analysis] 부분만 추출
        if "[Vision Analysis]" in vision_text:
            vision_part = vision_text.split("[Vision Analysis]")[1].split("[Original Content]")[0]
            print(vision_part)
        else:
            print(vision_text[:500])

    print("\n[정확도 체크]")

    # 표에서 숫자 추출
    actual_numbers = []
    for shape in slide_3.shapes:
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    import re
                    numbers = re.findall(r'\d+(?:,\d{3})*(?:\.\d+)?', cell.text)
                    actual_numbers.extend(numbers)

    print(f"실제 표의 숫자: {actual_numbers}")

    # Vision 분석에서 숫자 추출
    if 3 in slides_content and slides_content[3]["slide_summary"]:
        vision_text = slides_content[3]["slide_summary"]
        import re
        vision_numbers = re.findall(r'\d+(?:,\d{3})*(?:\.\d+)?', vision_text)
        print(f"Vision 분석의 숫자: {vision_numbers}")

        # 매칭 체크
        matched = 0
        for num in actual_numbers:
            if num in vision_numbers or num.replace(',', '') in [v.replace(',', '') for v in vision_numbers]:
                matched += 1

        accuracy = matched / len(actual_numbers) * 100 if actual_numbers else 100
        print(f"\n매칭률: {matched}/{len(actual_numbers)} ({accuracy:.1f}%)")

        if accuracy < 100:
            print("\n[누락된 숫자]")
            for num in actual_numbers:
                if num not in vision_numbers:
                    print(f"  - {num}")


if __name__ == "__main__":
    analyze_ppt_detailed()
