"""
PPT 파일 구조 진단: COM vs python-pptx 비교
"""

import sys
import os
from pathlib import Path

def diagnose_ppt_file():
    """PPT 파일을 COM과 python-pptx 두 가지 방식으로 분석하여 비교"""

    test_file = Path("data/test_pptx/advanced_01_financial_report.pptx")
    abs_path = str(test_file.absolute())

    print("=" * 80)
    print(f"PPT 파일 구조 진단: {test_file.name}")
    print("=" * 80)
    print()

    # 1. COM 방식으로 분석
    print("[1] Windows COM (PowerPoint) 분석")
    print("-" * 80)

    try:
        import win32com.client
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")

        try:
            presentation = powerpoint.Presentations.Open(abs_path)

            print(f"총 슬라이드 수: {presentation.Slides.Count}")
            print()

            for i in range(1, presentation.Slides.Count + 1):
                slide = presentation.Slides[i]

                # 제목 추출
                title = ""
                try:
                    if slide.Shapes.HasTitle:
                        title = slide.Shapes.Title.TextFrame.TextRange.Text
                except:
                    pass

                # 슬라이드 정보
                print(f"COM Slide {i}:")
                print(f"  SlideID: {slide.SlideID}")
                print(f"  SlideIndex: {slide.SlideIndex}")
                print(f"  제목: {title if title else '(제목 없음)'}")
                print(f"  Shape 개수: {slide.Shapes.Count}")

                # 숨김 여부
                try:
                    # SlideShowTransition 확인
                    hidden = slide.SlideShowTransition.Hidden
                    if hidden:
                        print(f"  [WARN] 숨겨진 슬라이드!")
                except:
                    pass

                print()

            presentation.Close()

        finally:
            powerpoint.Quit()

    except Exception as e:
        print(f"[ERROR] COM 분석 실패: {e}")

    print()

    # 2. python-pptx 방식으로 분석
    print("[2] python-pptx 분석")
    print("-" * 80)

    try:
        from pptx import Presentation

        prs = Presentation(abs_path)

        print(f"총 슬라이드 수: {len(prs.slides)}")
        print()

        for i, slide in enumerate(prs.slides, 1):
            # 제목 추출
            title = ""
            try:
                if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
                    title = slide.shapes.title.text.strip()
            except:
                pass

            print(f"python-pptx Slide {i}:")
            print(f"  제목: {title if title else '(제목 없음)'}")
            print(f"  Shape 개수: {len(slide.shapes)}")

            # 표 확인
            table_count = 0
            for shape in slide.shapes:
                if shape.has_table:
                    table_count += 1

            if table_count > 0:
                print(f"  표: {table_count}개")

            print()

    except Exception as e:
        print(f"[ERROR] python-pptx 분석 실패: {e}")

    print()

    # 3. 비교 분석
    print("=" * 80)
    print("[3] 비교 분석")
    print("=" * 80)
    print()
    print("위 두 결과를 비교하여 다음을 확인하세요:")
    print()
    print("1. 슬라이드 개수가 같은가?")
    print("2. 각 슬라이드의 제목이 같은 순서로 나타나는가?")
    print("3. COM에 숨겨진 슬라이드가 있는가?")
    print("4. SlideIndex와 순서가 일치하는가?")
    print()
    print("만약 불일치한다면:")
    print("  - COM이 보는 슬라이드 순서와 python-pptx가 보는 순서가 다름")
    print("  - 이미지-슬라이드 매칭 시 제목 기반 매칭 필요")
    print()


if __name__ == "__main__":
    diagnose_ppt_file()
