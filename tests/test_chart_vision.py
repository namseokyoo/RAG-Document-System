"""
Phase 1 차트 Vision 테스트

차트가 포함된 PPT 파일로 Vision 청킹 기능 검증
"""
import os
from pathlib import Path
from pptx import Presentation


def test_chart_vision():
    """차트 포함 PPT 파일로 Vision 청킹 테스트"""

    print("=" * 80)
    print("Phase 1: 차트 Vision 테스트")
    print("=" * 80)
    print()

    # 차트 테스트 파일
    test_file = Path("data/test_pptx/chart_test.pptx")

    if not test_file.exists():
        print(f"[ERROR] 테스트 파일 없음: {test_file}")
        print("먼저 create_chart_test_ppt.py를 실행하세요")
        return

    # 환경 설정
    from utils.pptx_chunking_engine import PPTXChunkingEngine
    from config import ConfigManager

    config_mgr = ConfigManager()
    config = config_mgr.get_all()

    llm_api_type = config.get("llm_api_type", "openai")
    llm_model = config.get("llm_model", "gpt-4o-mini")
    llm_api_key = config.get("llm_api_key", os.getenv("OPENAI_API_KEY"))

    print(f"설정:")
    print(f"  API 타입: {llm_api_type}")
    print(f"  모델: {llm_model}")
    print(f"  Vision 활성화: {config.get('enable_vision_chunking', True)}")
    print()

    # PPT 파일 분석
    prs = Presentation(str(test_file))
    slide_count = len(prs.slides)

    print(f"테스트 파일: {test_file.name}")
    print(f"총 슬라이드 수: {slide_count}")
    print()

    # 각 슬라이드 차트 확인
    print("슬라이드 구조:")
    print("-" * 80)

    chart_slides = []
    for i, slide in enumerate(prs.slides, 1):
        # 제목 추출
        title = ""
        try:
            if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
                title = slide.shapes.title.text.strip()
        except:
            pass

        if not title:
            title = f"제목없음-슬라이드{i}"

        # 차트 개수
        chart_count = 0
        table_count = 0
        for shape in slide.shapes:
            if shape.has_chart:
                chart_count += 1
                chart_slides.append(i)
            elif shape.has_table:
                table_count += 1

        print(f"  슬라이드 {i}: {title}")
        if chart_count > 0:
            print(f"    차트: {chart_count}개")
        if table_count > 0:
            print(f"    표: {table_count}개")

    print()
    print(f"차트 포함 슬라이드: {chart_slides}")
    print()

    # Vision 청킹 실행
    print("Vision 청킹 실행 중...")
    print("-" * 80)

    engine = PPTXChunkingEngine(config)

    chunks = engine.process_pptx_document(
        pptx_path=str(test_file.absolute()),
        enable_vision=True,
        llm_api_type=llm_api_type,
        llm_base_url=config.get('llm_base_url', ''),
        llm_model=llm_model,
        llm_api_key=llm_api_key
    )

    if not chunks:
        print("[ERROR] 청킹 실패 - 청크 0개")
        return

    print(f"총 {len(chunks)}개 청크 생성")
    print()

    # 결과 분석
    print("차트 분석 결과:")
    print("-" * 80)

    slide_summary_chunks = [c for c in chunks if c.chunk_type == "slide_summary"]

    success_count = 0
    fail_count = 0

    for slide_num in chart_slides:
        # 해당 슬라이드의 청크 찾기
        slide_chunks = [c for c in slide_summary_chunks if c.metadata.slide_number == slide_num]

        if not slide_chunks:
            print(f"\n슬라이드 {slide_num}: [SKIP] 청크 없음")
            continue

        chunk = slide_chunks[0]

        # 차트 분석 키워드 확인
        chart_keywords = ["차트", "막대", "선", "파이", "트렌드", "데이터 유형", "비교"]
        found_keywords = [kw for kw in chart_keywords if kw in chunk.content]

        is_chart_analyzed = len(found_keywords) >= 2  # 최소 2개 키워드

        if is_chart_analyzed:
            success_count += 1
            print(f"\n슬라이드 {slide_num}: [SUCCESS] 차트 분석 성공")
            print(f"  제목: {chunk.metadata.slide_title}")
            print(f"  발견 키워드: {', '.join(found_keywords)}")
            print(f"  내용 (앞 300자):")
            print(f"  {chunk.content[:300]}...")
        else:
            fail_count += 1
            print(f"\n슬라이드 {slide_num}: [FAIL] 차트 분석 실패")
            print(f"  제목: {chunk.metadata.slide_title}")
            print(f"  발견 키워드: {', '.join(found_keywords) if found_keywords else '없음'}")
            print(f"  내용 (앞 200자):")
            print(f"  {chunk.content[:200]}...")

    print()
    print("=" * 80)
    print("테스트 결과 요약")
    print("=" * 80)
    print(f"차트 슬라이드: {len(chart_slides)}개")
    print(f"분석 성공: {success_count}개")
    print(f"분석 실패: {fail_count}개")
    print(f"성공률: {success_count}/{len(chart_slides)} = {success_count/len(chart_slides)*100:.1f}%")
    print()

    # 목표: 90% 이상
    if success_count / len(chart_slides) >= 0.9:
        print("[SUCCESS] Phase 1 차트 Vision 테스트 통과 (90%+ 성공률)")
    else:
        print("[WARNING] Phase 1 차트 Vision 테스트 미달 (목표: 90%+)")

    print("=" * 80)


if __name__ == "__main__":
    try:
        test_chart_vision()
    except Exception as e:
        print(f"[ERROR] 테스트 실패: {e}")
        import traceback
        traceback.print_exc()
