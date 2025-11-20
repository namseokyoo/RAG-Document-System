"""
Phase 1-B 테스트: 표 구조 기반 이미지 매칭

이 스크립트는 제목 없는 슬라이드를 표 구조로 매칭하는 기능을 테스트합니다.
"""

import sys
import os
from pathlib import Path

def test_table_structure_matching():
    """표 구조 기반 매칭 테스트"""

    print("=" * 80)
    print("Phase 1-B: 표 구조 기반 이미지 매칭 테스트")
    print("=" * 80)
    print()

    # 환경 설정
    from utils.pptx_chunking_engine import PPTXChunkingEngine
    from config import ConfigManager

    # 테스트 파일
    test_file = Path("data/test_pptx/advanced_01_financial_report.pptx")

    if not test_file.exists():
        print(f"[ERROR] 테스트 파일을 찾을 수 없습니다: {test_file}")
        return

    print(f"테스트 파일: {test_file.name}")
    print()

    # config.json에서 설정 로드
    config_mgr = ConfigManager()
    config = config_mgr.get_all()

    llm_api_type = config.get("llm_api_type", "openai")
    llm_model = config.get("llm_model", "gpt-4o-mini")
    llm_api_key = config.get("llm_api_key", os.getenv("OPENAI_API_KEY"))

    print(f"설정:")
    print(f"  API 타입: {llm_api_type}")
    print(f"  모델: {llm_model}")
    print(f"  Vision 활성화: {config.get('enable_vision_chunking', False)}")
    print()

    # 청킹 엔진 초기화
    engine = PPTXChunkingEngine(config)

    # PPTX 파일 로드
    from pptx import Presentation

    prs = Presentation(str(test_file))
    print(f"총 슬라이드 수: {len(prs.slides)}")
    print()

    # 각 슬라이드의 표 구조 확인
    print("-" * 80)
    print("슬라이드별 표 구조 분석")
    print("-" * 80)
    print()

    for i, slide in enumerate(prs.slides, 1):
        # 제목 추출
        title = ""
        try:
            if slide.shapes.title and hasattr(slide.shapes.title, 'text'):
                title = slide.shapes.title.text.strip()
        except:
            pass

        if not title:
            title = f"제목없음-슬라이드{i}"

        print(f"슬라이드 {i}: {title}")

        # 표 구조 추출
        table_structure = engine._extract_table_structure(slide)

        if table_structure["has_table"]:
            print(f"  표 개수: {table_structure['table_count']}")
            for j, tbl in enumerate(table_structure["tables"], 1):
                print(f"    표{j}: {tbl['rows']}행 x {tbl['cols']}열")
                if tbl["headers"]:
                    print(f"      헤더: {', '.join(tbl['headers'][:3])}...")
        else:
            print(f"  표 없음")

        print()

    # Vision API로 COM 이미지의 표 구조 감지 테스트
    print("=" * 80)
    print("COM 렌더링 이미지 표 구조 감지 테스트")
    print("=" * 80)
    print()

    # COM으로 슬라이드 렌더링
    print("COM으로 슬라이드 렌더링 중...")

    # pptx_path 설정 (엔진이 필요로 함)
    engine.pptx_path = str(test_file.absolute())

    slide_images = engine._render_all_slides_via_com(len(prs.slides))

    if not slide_images:
        print("[ERROR] COM 렌더링 실패")
        return

    print(f"렌더링 완료: {len(slide_images)}개 이미지")
    print()

    # 각 이미지의 표 구조 감지
    for img_idx, img_data in slide_images.items():
        img_title = img_data.get("title", "")
        if not img_title:
            img_title = f"제목없음-슬라이드{img_idx + 1}"

        print(f"COM 이미지 {img_idx + 1}: {img_title}")

        # Vision API로 표 구조 감지
        print(f"  Vision API로 표 구조 감지 중...")
        img_structure = engine._detect_table_structure_via_vision(
            img_data["image"], llm_api_type, llm_api_key, llm_model
        )

        if img_structure["has_table"]:
            print(f"  표 개수: {img_structure['table_count']}")
            for j, tbl in enumerate(img_structure["tables"], 1):
                print(f"    표{j}: {tbl['rows']}행 x {tbl['cols']}열")
        else:
            print(f"  표 없음")

        print()

    # 매칭 테스트
    print("=" * 80)
    print("표 구조 기반 매칭 테스트")
    print("=" * 80)
    print()

    # 슬라이드 2 (제목 없음)를 매칭해보기
    slide_2 = prs.slides[1]  # 0-indexed

    print("테스트 대상: 슬라이드 2 (제목 없음)")
    print()

    matched_image = engine._match_by_table_structure(
        slide_2, 1, slide_images,  # slide_index = 1 (0-indexed)
        llm_api_type, llm_api_key, llm_model
    )

    if matched_image:
        print()
        print("[SUCCESS] 표 구조 기반 매칭 성공!")
    else:
        print()
        print("[FAIL] 표 구조 기반 매칭 실패")

    print()
    print("=" * 80)
    print("테스트 완료")
    print("=" * 80)


if __name__ == "__main__":
    try:
        test_table_structure_matching()
    except Exception as e:
        print(f"[ERROR] 테스트 실패: {e}")
        import traceback
        traceback.print_exc()
