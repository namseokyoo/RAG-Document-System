"""
Vision 청킹 정확도 평가 테스트
- 실제 PPT 내용 vs Vision 분석 결과 비교
- 표 숫자 정확도, 그래프 트렌드 인식, 맥락 보존 평가
"""

import sys
import time
from pathlib import Path
from typing import List, Dict, Any
import re

sys.path.insert(0, str(Path(__file__).parent))

from utils.pptx_chunking_engine import PPTXChunkingEngine
from config import ConfigManager
from pptx import Presentation


class VisionAccuracyEvaluator:
    """Vision 청킹 정확도 평가 클래스"""

    def __init__(self):
        self.evaluation_criteria = {
            "table_number_accuracy": {
                "name": "표 숫자 정확도",
                "description": "표에 있는 숫자가 Vision 분석에 포함되었는가",
                "weight": 0.4,
                "score": 0.0
            },
            "trend_recognition": {
                "name": "그래프 트렌드 인식",
                "description": "증가/감소/패턴을 올바르게 파악했는가",
                "weight": 0.2,
                "score": 0.0
            },
            "context_preservation": {
                "name": "맥락 보존",
                "description": "슬라이드의 핵심 메시지를 파악했는가",
                "weight": 0.2,
                "score": 0.0
            },
            "searchability": {
                "name": "검색 가능성",
                "description": "RAG 검색 시 찾을 수 있는 키워드가 포함되었는가",
                "weight": 0.2,
                "score": 0.0
            }
        }

    def extract_numbers_from_text(self, text: str) -> List[str]:
        """텍스트에서 숫자 추출 (억원, %, 등 포함)"""
        # 숫자 패턴: 1,234, 12.5%, 150억원 등
        patterns = [
            r'\d+(?:,\d{3})*(?:\.\d+)?%',  # 12.5%, 1,234%
            r'\d+(?:,\d{3})*(?:\.\d+)?억원?',  # 150억원, 1,234억
            r'\d+(?:,\d{3})*(?:\.\d+)?만원?',  # 500만원
            r'\d+(?:,\d{3})*(?:\.\d+)?원',  # 1,234원
            r'\d+(?:,\d{3})*(?:\.\d+)?',  # 일반 숫자
        ]

        numbers = []
        for pattern in patterns:
            matches = re.findall(pattern, text)
            numbers.extend(matches)

        return list(set(numbers))  # 중복 제거

    def extract_ground_truth(self, pptx_path: str) -> Dict[int, Dict[str, Any]]:
        """PPT 파일에서 실제 내용(Ground Truth) 추출"""
        prs = Presentation(pptx_path)
        ground_truth = {}

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_data = {
                "title": "",
                "text": "",
                "numbers": [],
                "has_table": False,
                "has_chart": False,
                "table_data": [],
                "bullet_points": []
            }

            # 제목 추출
            if slide.shapes.title:
                slide_data["title"] = slide.shapes.title.text

            # 모든 텍스트와 구조 추출
            full_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    full_text.append(shape.text)

                # 표 확인
                if shape.has_table:
                    slide_data["has_table"] = True
                    table = shape.table
                    for row in table.rows:
                        row_data = [cell.text for cell in row.cells]
                        slide_data["table_data"].append(row_data)

                # 차트 확인
                if shape.has_chart:
                    slide_data["has_chart"] = True

                # 불릿 포인트 추출
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_data["bullet_points"].append(paragraph.text.strip())

            slide_data["text"] = "\n".join(full_text)
            slide_data["numbers"] = self.extract_numbers_from_text(slide_data["text"])

            ground_truth[slide_idx] = slide_data

        return ground_truth

    def evaluate_table_accuracy(self, ground_truth: List[str], vision_text: str) -> float:
        """표 숫자 정확도 평가"""
        if not ground_truth:
            return 1.0  # 숫자가 없으면 만점

        # Vision 분석에서 추출된 숫자
        vision_numbers = self.extract_numbers_from_text(vision_text)

        # 일치하는 숫자 개수
        matched = 0
        for gt_num in ground_truth:
            # 정확한 일치 또는 유사한 형태 (쉼표 제거 등)
            gt_clean = gt_num.replace(',', '')
            for v_num in vision_numbers:
                v_clean = v_num.replace(',', '')
                if gt_clean == v_clean or gt_clean in v_clean or v_clean in gt_clean:
                    matched += 1
                    break

        accuracy = matched / len(ground_truth) if ground_truth else 0.0
        return accuracy

    def evaluate_trend_recognition(self, ground_truth_data: Dict, vision_text: str) -> float:
        """트렌드 인식 평가"""
        score = 0.0
        checks = 0

        # 증가/감소 키워드 확인
        trend_keywords = ['증가', '감소', '성장', '하락', '상승', '하향', '+', '-', '%']

        if any(keyword in ground_truth_data.get("text", "") for keyword in trend_keywords):
            checks += 1
            if any(keyword in vision_text for keyword in trend_keywords):
                score += 1

        # 비교 표현 확인
        comparison_keywords = ['대비', '전년', '전기', 'vs', 'YoY', 'QoQ']
        if any(keyword in ground_truth_data.get("text", "") for keyword in comparison_keywords):
            checks += 1
            if any(keyword in vision_text for keyword in comparison_keywords):
                score += 1

        return score / checks if checks > 0 else 1.0

    def evaluate_context_preservation(self, ground_truth_data: Dict, vision_text: str) -> float:
        """맥락 보존 평가"""
        score = 0.0
        checks = 0

        # 제목이 Vision 분석에 포함되어 있는가
        title = ground_truth_data.get("title", "")
        if title:
            checks += 1
            # 제목의 핵심 단어 추출 (조사 제외)
            title_words = [w for w in title.split() if len(w) > 1]
            if any(word in vision_text for word in title_words):
                score += 1

        # 주요 불릿 포인트가 포함되어 있는가
        bullet_points = ground_truth_data.get("bullet_points", [])
        if bullet_points:
            checks += len(bullet_points[:3])  # 상위 3개만 체크
            for bullet in bullet_points[:3]:
                bullet_words = [w for w in bullet.split() if len(w) > 1]
                if any(word in vision_text for word in bullet_words[:3]):  # 각 불릿의 첫 3단어
                    score += 1

        return score / checks if checks > 0 else 1.0

    def evaluate_searchability(self, ground_truth_data: Dict, vision_text: str) -> float:
        """검색 가능성 평가 - RAG에서 찾을 수 있는 키워드"""
        score = 0.0
        checks = 0

        # 중요 키워드 (명사, 숫자)
        important_keywords = []

        # 숫자는 검색의 핵심
        gt_numbers = ground_truth_data.get("numbers", [])
        if gt_numbers:
            checks += min(len(gt_numbers), 5)  # 최대 5개까지
            for num in gt_numbers[:5]:
                if num in vision_text:
                    score += 1

        # 제목의 핵심 단어
        title = ground_truth_data.get("title", "")
        if title:
            title_keywords = [w for w in title.split() if len(w) > 2]
            checks += min(len(title_keywords), 3)
            for kw in title_keywords[:3]:
                if kw in vision_text:
                    score += 1

        return score / checks if checks > 0 else 1.0

    def evaluate_slide(self, slide_num: int, ground_truth: Dict, vision_text: str) -> Dict[str, float]:
        """개별 슬라이드 평가"""
        scores = {
            "table_number_accuracy": self.evaluate_table_accuracy(
                ground_truth.get("numbers", []),
                vision_text
            ),
            "trend_recognition": self.evaluate_trend_recognition(
                ground_truth,
                vision_text
            ),
            "context_preservation": self.evaluate_context_preservation(
                ground_truth,
                vision_text
            ),
            "searchability": self.evaluate_searchability(
                ground_truth,
                vision_text
            )
        }

        return scores

    def calculate_overall_score(self, slide_scores: List[Dict[str, float]]) -> Dict[str, Any]:
        """전체 점수 계산"""
        if not slide_scores:
            return {}

        # 각 지표별 평균
        avg_scores = {}
        for criterion in self.evaluation_criteria.keys():
            criterion_scores = [s[criterion] for s in slide_scores if criterion in s]
            avg_scores[criterion] = sum(criterion_scores) / len(criterion_scores) if criterion_scores else 0.0

        # 가중 평균
        weighted_score = sum(
            avg_scores[criterion] * self.evaluation_criteria[criterion]["weight"]
            for criterion in self.evaluation_criteria.keys()
        )

        return {
            "criterion_scores": avg_scores,
            "weighted_total": weighted_score,
            "slide_count": len(slide_scores)
        }


def test_complex_ppt():
    """복잡한 PPT 파일로 정확도 테스트"""
    print("=" * 80)
    print("Vision 청킹 정확도 평가 테스트")
    print("=" * 80)

    # 테스트할 파일들
    test_files = [
        "data/test_pptx/advanced_01_financial_report.pptx",
        "data/test_pptx/complex_05_comprehensive_report.pptx",
        "data/test_pptx/complex_03_data_analysis_report.pptx"
    ]

    config_mgr = ConfigManager()
    config = config_mgr.get_all()

    evaluator = VisionAccuracyEvaluator()

    all_results = []

    for test_file_path in test_files:
        test_file = Path(test_file_path)

        if not test_file.exists():
            print(f"\n[SKIP] 파일 없음: {test_file}")
            continue

        print(f"\n" + "=" * 80)
        print(f"테스트 파일: {test_file.name}")
        print(f"파일 크기: {test_file.stat().st_size / 1024:.2f} KB")
        print("=" * 80)

        # Ground Truth 추출
        print("\n[1단계] 실제 내용(Ground Truth) 추출 중...")
        ground_truth = evaluator.extract_ground_truth(str(test_file))

        print(f"  - 총 슬라이드: {len(ground_truth)}")
        for slide_num, data in ground_truth.items():
            print(f"  - 슬라이드 {slide_num}: {data['title'][:50] if data['title'] else '(제목 없음)'}")
            print(f"    표: {data['has_table']}, 차트: {data['has_chart']}, 숫자: {len(data['numbers'])}개")

        # Vision 청킹 실행
        print("\n[2단계] Vision 청킹 실행 중...")
        engine = PPTXChunkingEngine(config)

        start_time = time.time()
        chunks = engine.process_pptx_document(
            pptx_path=str(test_file),
            enable_vision=True,
            llm_api_type=config['llm_api_type'],
            llm_base_url=config.get('llm_base_url', ''),
            llm_model=config['llm_model'],
            llm_api_key=config.get('llm_api_key', '')
        )
        elapsed_time = time.time() - start_time

        print(f"  - 처리 시간: {elapsed_time:.2f}초")
        print(f"  - 생성된 청크: {len(chunks)}개")

        # Vision 분석 텍스트만 추출
        vision_texts = {}
        for chunk in chunks:
            slide_num = chunk.metadata.slide_number
            chunk_text = chunk.content if hasattr(chunk, 'content') else chunk.text if hasattr(chunk, 'text') else str(chunk)

            if slide_num not in vision_texts:
                vision_texts[slide_num] = ""

            vision_texts[slide_num] += chunk_text + "\n"

        # 평가
        print("\n[3단계] 정확도 평가 중...")
        slide_scores = []

        for slide_num in ground_truth.keys():
            if slide_num not in vision_texts:
                print(f"  [WARN] 슬라이드 {slide_num}: Vision 분석 없음")
                continue

            scores = evaluator.evaluate_slide(
                slide_num,
                ground_truth[slide_num],
                vision_texts[slide_num]
            )
            slide_scores.append(scores)

            print(f"\n  슬라이드 {slide_num}: {ground_truth[slide_num]['title'][:40]}")
            print(f"    - 표 숫자 정확도: {scores['table_number_accuracy']*100:.1f}%")
            print(f"    - 트렌드 인식: {scores['trend_recognition']*100:.1f}%")
            print(f"    - 맥락 보존: {scores['context_preservation']*100:.1f}%")
            print(f"    - 검색 가능성: {scores['searchability']*100:.1f}%")

        # 전체 점수
        overall = evaluator.calculate_overall_score(slide_scores)

        print("\n" + "-" * 80)
        print(f"[종합 평가] {test_file.name}")
        print("-" * 80)
        print(f"표 숫자 정확도: {overall['criterion_scores']['table_number_accuracy']*100:.1f}%")
        print(f"트렌드 인식: {overall['criterion_scores']['trend_recognition']*100:.1f}%")
        print(f"맥락 보존: {overall['criterion_scores']['context_preservation']*100:.1f}%")
        print(f"검색 가능성: {overall['criterion_scores']['searchability']*100:.1f}%")
        print(f"\n가중 평균 점수: {overall['weighted_total']*100:.1f}점")

        # 등급 판정
        score = overall['weighted_total'] * 100
        if score >= 90:
            grade = "A (매우 우수)"
        elif score >= 80:
            grade = "B (우수)"
        elif score >= 70:
            grade = "C (양호)"
        elif score >= 60:
            grade = "D (보통)"
        else:
            grade = "F (개선 필요)"

        print(f"종합 등급: {grade}")

        all_results.append({
            "file": test_file.name,
            "overall_score": overall,
            "processing_time": elapsed_time
        })

    # 최종 요약
    print("\n" + "=" * 80)
    print("전체 테스트 요약")
    print("=" * 80)

    for result in all_results:
        print(f"\n{result['file']}: {result['overall_score']['weighted_total']*100:.1f}점")

    if all_results:
        avg_score = sum(r['overall_score']['weighted_total'] for r in all_results) / len(all_results)
        print(f"\n평균 점수: {avg_score*100:.1f}점")

        print("\n[결론]")
        if avg_score >= 0.8:
            print("Vision 청킹이 매우 효과적입니다. 프로덕션 사용 권장합니다.")
        elif avg_score >= 0.7:
            print("Vision 청킹이 효과적입니다. 일부 개선 후 프로덕션 사용 가능합니다.")
        elif avg_score >= 0.6:
            print("Vision 청킹이 보통 수준입니다. 추가 최적화가 필요합니다.")
        else:
            print("Vision 청킹의 정확도가 낮습니다. 프롬프트 개선이 필요합니다.")


if __name__ == "__main__":
    test_complex_ppt()
